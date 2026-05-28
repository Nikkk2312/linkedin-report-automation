"""Unit tests for src/report_generator.py"""

import json
import os
import tempfile
import unittest

from src.report_generator import generate_report, _safe_float


def _make_campaign(campaign_id, name, display_name):
    """Create a minimal campaign dict for testing."""
    return {
        "id": str(campaign_id),
        "name": name,
        "display_name": display_name,
        "campaign_type": "Engagement",
        "status": "ACTIVE",
        "start_date": "2025-12-08",
        "geo_display": "US, Canada",
        "impressions": 10000,
        "clicks": 200,
        "engagements": 500,
        "likes": 100,
        "reactions": 120,
        "comments": 30,
        "shares": 20,
        "follows": 10,
        "cost_usd": 150.0,
        "landing_page_clicks": 180,
        "company_page_clicks": 50,
        "other_engagements": 40,
        "conversions": 5,
        "one_click_leads": 2,
        "video_views": 3000,
        "video_completions": 1000,
        "video_first_quartile": 2500,
        "video_midpoint": 2000,
        "video_third_quartile": 1500,
        "full_screen_plays": 50,
        "ctr": 2.0,
        "cpc": 0.75,
        "cpm": 15.0,
        "cost_per_engagement": 0.30,
        "engagement_rate": 5.0,
        "monthly_trends": [
            {"month": "Jan 2026", "impressions": 5000, "clicks": 100, "cost": 75.0, "engagements": 250},
            {"month": "Feb 2026", "impressions": 5000, "clicks": 100, "cost": 75.0, "engagements": 250},
        ],
        "creatives": [],
        "demographics": {},
    }


def _make_test_data(num_campaigns):
    """Create test JSON data with the specified number of campaigns."""
    campaigns = []
    for i in range(num_campaigns):
        campaigns.append(_make_campaign(1000 + i, f"Campaign {i}", f"Camp {i}"))
    return {"report_date": "24-05-2026", "campaigns": campaigns}


class TestSafeFloat(unittest.TestCase):
    """Tests for _safe_float()."""

    def test_none(self):
        self.assertIsNone(_safe_float(None))

    def test_integer(self):
        self.assertEqual(_safe_float(42), 42.0)

    def test_float(self):
        self.assertEqual(_safe_float(3.14), 3.14)

    def test_string_number(self):
        self.assertEqual(_safe_float('123.45'), 123.45)

    def test_invalid_string(self):
        self.assertIsNone(_safe_float('abc'))

    def test_zero(self):
        self.assertEqual(_safe_float(0), 0.0)

    def test_negative(self):
        self.assertEqual(_safe_float(-5.5), -5.5)

    def test_empty_string(self):
        self.assertIsNone(_safe_float(''))


class TestGenerateReport(unittest.TestCase):
    """Tests for generate_report()."""

    def _generate_and_check(self, num_campaigns):
        """Helper: generate a report with N campaigns and return the PPTX path."""
        data = _make_test_data(num_campaigns)
        with tempfile.TemporaryDirectory() as tmpdir:
            json_path = os.path.join(tmpdir, 'test_data.json')
            pptx_path = os.path.join(tmpdir, 'test_report.pptx')
            with open(json_path, 'w', encoding='utf-8') as f:
                json.dump(data, f)
            generate_report(json_path, pptx_path, logo_path=None)
            self.assertTrue(os.path.isfile(pptx_path), "PPTX file should exist")
            # Verify the file is a valid PPTX by loading it
            from pptx import Presentation
            prs = Presentation(pptx_path)
            return prs

    def test_generate_report_single_campaign(self):
        """Generate report with 1 campaign, verify output exists and has slides."""
        prs = self._generate_and_check(1)
        self.assertGreater(len(prs.slides), 0, "Report should have at least 1 slide")

    def test_generate_report_two_campaigns(self):
        """Generate report with 2 campaigns, verify more slides than 1 campaign."""
        prs1 = self._generate_and_check(1)
        prs2 = self._generate_and_check(2)
        self.assertGreater(len(prs2.slides), len(prs1.slides),
                           "2-campaign report should have more slides than 1-campaign report")

    def test_generate_report_three_campaigns(self):
        """Generate report with 3 campaigns, verify more slides than 2 campaigns."""
        prs2 = self._generate_and_check(2)
        prs3 = self._generate_and_check(3)
        self.assertGreater(len(prs3.slides), len(prs2.slides),
                           "3-campaign report should have more slides than 2-campaign report")

    def test_generate_report_with_mock_data_file(self):
        """Generate report from existing test_data.json if available."""
        test_data_path = os.path.join(
            os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
            'output', 'test_data.json'
        )
        if not os.path.isfile(test_data_path):
            self.skipTest("test_data.json not found in output/")
        with tempfile.TemporaryDirectory() as tmpdir:
            pptx_path = os.path.join(tmpdir, 'test_report.pptx')
            generate_report(test_data_path, pptx_path, logo_path=None)
            self.assertTrue(os.path.isfile(pptx_path))
            from pptx import Presentation
            prs = Presentation(pptx_path)
            self.assertGreater(len(prs.slides), 5, "Full report should have many slides")


if __name__ == '__main__':
    unittest.main()
