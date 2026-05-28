"""Low-level PPTX helper functions for cell formatting, shapes, and layout."""

import os
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

from .config import (
    SLIDE_WIDTH, CLR_WHITE, CLR_LINKEDIN_BLUE, CLR_ORANGE, CLR_ACCENT_DARK,
    CLR_BORDER, FONT_NAME, LOGO_LEFT, LOGO_TOP, LOGO_WIDTH, LOGO_HEIGHT,
)


def set_cell_border(cell, color=CLR_BORDER, width=Emu(12700)):
    """Set borders on a table cell."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for border_name in ['lnL', 'lnR', 'lnT', 'lnB']:
        ln = tcPr.find(qn(f'a:{border_name}'))
        if ln is not None:
            tcPr.remove(ln)
        ln = tcPr.makeelement(qn(f'a:{border_name}'), {'w': str(int(width))})
        solid = ln.makeelement(qn('a:solidFill'), {})
        srgb = solid.makeelement(qn('a:srgbClr'), {'val': str(color)})
        solid.append(srgb)
        ln.append(solid)
        tcPr.append(ln)


def set_cell_fill(cell, color):
    """Set solid fill on a table cell."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for existing in tcPr.findall(qn('a:solidFill')):
        tcPr.remove(existing)
    for existing in tcPr.findall(qn('a:noFill')):
        tcPr.remove(existing)
    solid = tcPr.makeelement(qn('a:solidFill'), {})
    srgb = solid.makeelement(qn('a:srgbClr'), {'val': str(color)})
    solid.append(srgb)
    tcPr.insert(0, solid)


def set_cell_text(cell, text, font_size=Pt(13), bold=False, color=None,
                  alignment=PP_ALIGN.LEFT, font_name=FONT_NAME):
    """Set text in a table cell with formatting."""
    from .config import CLR_BLACK
    if color is None:
        color = CLR_BLACK
    cell.text = ''
    p = cell.text_frame.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = str(text)
    run.font.name = font_name
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_logo(slide, logo_path):
    """Add company logo to the top-right corner of a slide."""
    if logo_path and os.path.exists(logo_path):
        slide.shapes.add_picture(logo_path, LOGO_LEFT, LOGO_TOP, LOGO_WIDTH, LOGO_HEIGHT)


def add_slide_bg(slide):
    """Ensure white background."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = CLR_WHITE


def add_accent_bar(slide, top=None):
    """Add a thin accent bar across the bottom of a slide."""
    bar_top = top or Emu(6500000)
    shape = slide.shapes.add_shape(1, Emu(0), bar_top, SLIDE_WIDTH, Emu(80000))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CLR_LINKEDIN_BLUE
    shape.line.fill.background()


def add_top_line(slide):
    """Add a subtle separator line below the header area."""
    shape = slide.shapes.add_shape(
        1, Emu(274320), Emu(685800), Emu(8595360), Emu(19050)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CLR_LINKEDIN_BLUE
    shape.line.fill.background()


def add_top_band(slide):
    """Add blue accent band at the top of a slide."""
    band = slide.shapes.add_shape(1, Emu(0), Emu(0), SLIDE_WIDTH, Emu(1200000))
    band.fill.solid()
    band.fill.fore_color.rgb = CLR_LINKEDIN_BLUE
    band.line.fill.background()
    return band


def add_orange_accent(slide, left, top, width=Emu(2400000)):
    """Add a thin orange accent line."""
    accent = slide.shapes.add_shape(1, left, top, width, Emu(38100))
    accent.fill.solid()
    accent.fill.fore_color.rgb = CLR_ORANGE
    accent.line.fill.background()
    return accent


def add_heading(slide, text, left=Emu(274320), top=Emu(137160),
                width=Emu(8229600), height=Emu(548640)):
    """Add a consistent heading to a slide."""
    # Truncate long titles to prevent overlap with logo
    max_chars = 48
    display_text = text if len(text) <= max_chars else text[:max_chars - 1] + '\u2026'
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = display_text
    run.font.name = FONT_NAME
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = CLR_ACCENT_DARK
    return txBox


def add_hyperlink(slide, run, url):
    """Add a hyperlink to a text run."""
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
    r_elem = run._r
    rPr = r_elem.get_or_add_rPr()
    rel_id = slide.part.relate_to(url, RT.HYPERLINK, is_external=True)
    hlinkClick = rPr.makeelement(qn('a:hlinkClick'), {qn('r:id'): rel_id})
    rPr.append(hlinkClick)
