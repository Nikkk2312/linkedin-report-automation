"""
LinkedIn Ad Campaign Report Generator
Generates a professional PPTX report with executive summary, cost analysis,
engagement breakdown, monthly trends, creative rankings, and demographic insights.
"""

import json
import sys
import os
import math
import requests
from datetime import datetime
from io import BytesIO
from collections import defaultdict

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor

from .config import (
    SLIDE_WIDTH, SLIDE_HEIGHT, CLR_LINKEDIN_BLUE, CLR_ACCENT_DARK, CLR_ORANGE,
    CLR_BLACK, CLR_WHITE, CLR_ROW_LIGHT, CLR_BORDER, CLR_DARK_GRAY,
    CLR_SUBTLE_GRAY, CLR_LIGHT_GRAY, CLR_SUCCESS_GREEN, CLR_DANGER_RED,
    FONT_NAME, FONT_SIZE_TITLE, FONT_SIZE_HEADING, FONT_SIZE_BODY,
    FONT_SIZE_SMALL, FONT_SIZE_CAPTION, CTR_BENCHMARKS,
    LOGO_LEFT, LOGO_TOP, LOGO_WIDTH, LOGO_HEIGHT,
)
from .pptx_helpers import (
    set_cell_border, set_cell_fill, set_cell_text, add_logo, add_slide_bg,
    add_accent_bar, add_top_line, add_top_band, add_orange_accent,
    add_heading, add_hyperlink,
)
from .formatters import (
    format_number, format_ctr, calc_ctr_value, format_date_ordinal,
    abbreviate_number, format_currency, format_percentage,
)


def download_image(url, timeout=10):
    """Download an image and return BytesIO."""
    try:
        r = requests.get(url, timeout=timeout)
        if r.status_code == 200:
            return BytesIO(r.content)
    except Exception as e:
        print(f"  Warning: Could not download image: {e}", file=sys.stderr)
    return None


def _safe_float(val):
    """Safely convert to float, returning None on failure."""
    if val is None:
        return None
    try:
        return float(val)
    except (ValueError, TypeError):
        return None


def _add_kpi_card(slide, left, top, width, height, value_text, label_text,
                  subtitle_text='', value_color=None):
    """Add a KPI card shape with value, label, and optional subtitle."""
    if value_color is None:
        value_color = CLR_ACCENT_DARK

    card = slide.shapes.add_shape(1, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = CLR_LIGHT_GRAY
    card.line.color.rgb = CLR_BORDER
    card.line.width = Emu(12700)

    # Value
    txVal = slide.shapes.add_textbox(left + Emu(80000), top + Emu(200000),
                                      width - Emu(160000), Emu(450000))
    pv = txVal.text_frame.paragraphs[0]
    pv.alignment = PP_ALIGN.CENTER
    rv = pv.add_run()
    rv.text = value_text
    rv.font.name = FONT_NAME
    rv.font.size = Pt(28)
    rv.font.bold = True
    rv.font.color.rgb = value_color

    # Label
    txLbl = slide.shapes.add_textbox(left + Emu(80000), top + Emu(680000),
                                      width - Emu(160000), Emu(280000))
    pl = txLbl.text_frame.paragraphs[0]
    pl.alignment = PP_ALIGN.CENTER
    rl = pl.add_run()
    rl.text = label_text
    rl.font.name = FONT_NAME
    rl.font.size = Pt(10)
    rl.font.bold = True
    rl.font.color.rgb = CLR_DARK_GRAY

    if subtitle_text:
        txSub = slide.shapes.add_textbox(left + Emu(80000), top + Emu(960000),
                                          width - Emu(160000), Emu(250000))
        ps = txSub.text_frame.paragraphs[0]
        ps.alignment = PP_ALIGN.CENTER
        rs = ps.add_run()
        rs.text = subtitle_text
        rs.font.name = FONT_NAME
        rs.font.size = FONT_SIZE_CAPTION
        rs.font.color.rgb = CLR_SUBTLE_GRAY


def _add_slide_number(slide, slide_num, total_slides):
    """Add 'Slide X of Y' text to the bottom-right corner of a slide."""
    tx = slide.shapes.add_textbox(
        Emu(7200000), Emu(6400000), Emu(1800000), Emu(300000)
    )
    p = tx.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = f'Slide {slide_num} of {total_slides}'
    r.font.name = FONT_NAME
    r.font.size = FONT_SIZE_CAPTION
    r.font.color.rgb = CLR_SUBTLE_GRAY


# --- SLIDE GENERATORS --------------------------------------------------------


def create_title_slide(prs, report_date, logo_path):
    """Slide 1: Professional title slide with blue band."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)
    add_top_band(slide)

    txBrand = slide.shapes.add_textbox(Emu(457200), Emu(350000), Emu(5000000), Emu(500000))
    pb = txBrand.text_frame.paragraphs[0]
    rb = pb.add_run()
    rb.text = 'LinkedIn'
    rb.font.name = FONT_NAME
    rb.font.size = Pt(28)
    rb.font.bold = True
    rb.font.color.rgb = CLR_WHITE

    txBox = slide.shapes.add_textbox(Emu(457200), Emu(2100000), Emu(8229600), Emu(914400))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = 'Ad Campaign Performance Report'
    run.font.name = FONT_NAME
    run.font.size = FONT_SIZE_TITLE
    run.font.bold = True
    run.font.color.rgb = CLR_ACCENT_DARK

    add_orange_accent(slide, Emu(457200), Emu(2980000))

    txBox2 = slide.shapes.add_textbox(Emu(457200), Emu(3050000), Emu(8229600), Emu(731520))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.LEFT
    run2 = p2.add_run()
    run2.text = f'Report Date: {report_date}'
    run2.font.name = FONT_NAME
    run2.font.size = Pt(20)
    run2.font.color.rgb = CLR_SUBTLE_GRAY

    add_accent_bar(slide)
    add_logo(slide, logo_path)


def create_table_of_contents_slide(prs, campaigns, logo_path):
    """Table of Contents slide listing all report sections."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)
    add_heading(slide, 'Table of Contents')
    add_top_line(slide)

    txBox = slide.shapes.add_textbox(Emu(500000), Emu(900000), Emu(8100000), Emu(5000000))
    tf = txBox.text_frame
    tf.word_wrap = True

    sections = [
        'Executive Summary',
        'Key Insights',
        'Campaign Overview',
        'Impressions Comparison',
        'Campaign Ranking',
        'CTR Comparison',
        'Engagement Breakdown',
        'Cost & Efficiency Analysis',
    ]

    for i, section in enumerate(sections):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(4)
        r = p.add_run()
        r.text = f'  {section}'
        r.font.name = FONT_NAME
        r.font.size = Pt(13)
        r.font.color.rgb = CLR_DARK_GRAY
        r.font.bold = True

    # Add per-campaign sections
    p_header = tf.add_paragraph()
    p_header.space_before = Pt(10)
    p_header.space_after = Pt(4)
    rh = p_header.add_run()
    rh.text = '  Per-Campaign Deep Dives:'
    rh.font.name = FONT_NAME
    rh.font.size = Pt(13)
    rh.font.bold = True
    rh.font.color.rgb = CLR_LINKEDIN_BLUE

    for camp in campaigns:
        name = camp.get('display_name', camp.get('name', 'Unknown'))
        p_camp = tf.add_paragraph()
        p_camp.space_after = Pt(2)
        rc = p_camp.add_run()
        rc.text = f'      {name}'
        rc.font.name = FONT_NAME
        rc.font.size = Pt(12)
        rc.font.color.rgb = CLR_DARK_GRAY

    # Top Creatives + Thank You
    for section in ['Top Creatives Across All Campaigns', 'Thank You']:
        p_s = tf.add_paragraph()
        p_s.space_after = Pt(4)
        rs = p_s.add_run()
        rs.text = f'  {section}'
        rs.font.name = FONT_NAME
        rs.font.size = Pt(13)
        rs.font.color.rgb = CLR_DARK_GRAY
        rs.font.bold = True

    add_accent_bar(slide)
    add_logo(slide, logo_path)


def create_executive_summary_slide(prs, campaigns, logo_path):
    """Slide 2: Executive summary with 5 KPI cards + top performer."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)
    add_heading(slide, 'Executive Summary')
    add_top_line(slide)

    total_impressions = sum(c.get('impressions') or 0 for c in campaigns)
    total_clicks = sum(c.get('clicks') or 0 for c in campaigns)
    total_engagements = sum(c.get('engagements') or 0 for c in campaigns)
    total_spend = sum(_safe_float(c.get('cost_usd')) or 0 for c in campaigns)
    active_count = sum(1 for c in campaigns if c.get('status') == 'ACTIVE')
    overall_ctr = calc_ctr_value(total_impressions, total_clicks)

    top_campaign = max(campaigns, key=lambda c: c.get('impressions') or 0)

    # 5 KPI cards
    kpis = [
        ('Total Impressions', abbreviate_number(total_impressions), f'{len(campaigns)} campaigns'),
        ('Total Clicks', abbreviate_number(total_clicks), f'{active_count} active'),
        ('Overall CTR', f'{overall_ctr:.2f}%' if overall_ctr else 'NA', 'click-through rate'),
        ('Total Engagements', abbreviate_number(total_engagements), 'likes, shares, comments'),
        ('Total Spend', format_currency(total_spend) if total_spend > 0 else 'NA', 'USD'),
    ]

    card_width = Emu(1600000)
    card_height = Emu(1350000)
    card_gap = Emu(100000)
    start_left = Emu(250000)
    card_top = Emu(900000)

    for i, (label, value, subtitle) in enumerate(kpis):
        left = start_left + i * (card_width + card_gap)
        _add_kpi_card(slide, left, card_top, card_width, card_height,
                      value, label, subtitle)

    # Top performer callout
    callout_top = Emu(2550000)
    txCallout = slide.shapes.add_textbox(Emu(250000), callout_top, Emu(8600000), Emu(500000))
    tf = txCallout.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = 'Top Performer:  '
    r1.font.name = FONT_NAME
    r1.font.size = Pt(13)
    r1.font.bold = True
    r1.font.color.rgb = CLR_LINKEDIN_BLUE

    r2 = p.add_run()
    top_name = top_campaign.get('display_name', top_campaign.get('name', 'Unknown'))
    top_imp = format_number(top_campaign.get('impressions'))
    top_ctr_val = calc_ctr_value(top_campaign.get('impressions'), top_campaign.get('clicks'))
    top_ctr = f'{top_ctr_val:.2f}%' if top_ctr_val else 'NA'
    r2.text = f'{top_name}  |  {top_imp} impressions  |  {top_ctr} CTR'
    r2.font.name = FONT_NAME
    r2.font.size = Pt(13)
    r2.font.color.rgb = CLR_DARK_GRAY

    # Efficiency metrics row (if spend data exists)
    if total_spend > 0:
        overall_cpc = total_spend / total_clicks if total_clicks else None
        overall_cpm = (total_spend / total_impressions * 1000) if total_impressions else None
        overall_cpe = total_spend / total_engagements if total_engagements else None

        eff_top = Emu(3100000)
        eff_metrics = [
            ('Avg. CPC', format_currency(overall_cpc) if overall_cpc else 'NA'),
            ('Avg. CPM', format_currency(overall_cpm) if overall_cpm else 'NA'),
            ('Cost/Engagement', format_currency(overall_cpe) if overall_cpe else 'NA'),
        ]

        eff_card_w = Emu(2700000)
        eff_card_h = Emu(1000000)
        eff_gap = Emu(120000)
        eff_start = Emu(350000)

        for i, (label, value) in enumerate(eff_metrics):
            left = eff_start + i * (eff_card_w + eff_gap)
            _add_kpi_card(slide, left, eff_top, eff_card_w, eff_card_h,
                          value, label)

    # "vs" comparison for exactly 2 campaigns
    if len(campaigns) == 2:
        c1 = campaigns[0]
        c2 = campaigns[1]
        c1_name = c1.get('display_name', c1.get('name', 'Campaign 1'))[:18]
        c2_name = c2.get('display_name', c2.get('name', 'Campaign 2'))[:18]

        comparisons = []
        # CTR comparison
        ctr1 = calc_ctr_value(c1.get('impressions'), c1.get('clicks'))
        ctr2 = calc_ctr_value(c2.get('impressions'), c2.get('clicks'))
        if ctr1 is not None and ctr2 is not None:
            winner = c1_name if ctr1 >= ctr2 else c2_name
            comparisons.append(f'CTR: {winner}')

        # Impressions comparison
        imp1 = c1.get('impressions') or 0
        imp2 = c2.get('impressions') or 0
        winner = c1_name if imp1 >= imp2 else c2_name
        comparisons.append(f'Impressions: {winner}')

        # Engagement rate comparison
        er1 = _safe_float(c1.get('engagement_rate'))
        er2 = _safe_float(c2.get('engagement_rate'))
        if er1 is not None and er2 is not None:
            winner = c1_name if er1 >= er2 else c2_name
            comparisons.append(f'Eng. Rate: {winner}')

        # CPC comparison (lower is better)
        cpc1 = _safe_float(c1.get('cpc'))
        cpc2 = _safe_float(c2.get('cpc'))
        if cpc1 is not None and cpc2 is not None:
            winner = c1_name if cpc1 <= cpc2 else c2_name
            comparisons.append(f'CPC (lower): {winner}')

        if comparisons:
            vs_top = Emu(4300000) if total_spend > 0 else Emu(3200000)
            txVs = slide.shapes.add_textbox(Emu(250000), vs_top, Emu(8600000), Emu(500000))
            tf_vs = txVs.text_frame
            tf_vs.word_wrap = True
            p_vs = tf_vs.paragraphs[0]
            rv1 = p_vs.add_run()
            rv1.text = f'{c1_name} vs {c2_name}:  '
            rv1.font.name = FONT_NAME
            rv1.font.size = Pt(11)
            rv1.font.bold = True
            rv1.font.color.rgb = CLR_ORANGE
            rv2 = p_vs.add_run()
            rv2.text = '  |  '.join(comparisons)
            rv2.font.name = FONT_NAME
            rv2.font.size = Pt(11)
            rv2.font.color.rgb = CLR_DARK_GRAY

    add_accent_bar(slide)
    add_logo(slide, logo_path)


def create_key_insights_slide(prs, campaigns, logo_path):
    """Key Insights slide with auto-generated data-driven observations."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)
    add_heading(slide, 'Key Insights')
    add_top_line(slide)

    total_impressions = sum(c.get('impressions') or 0 for c in campaigns)
    total_clicks = sum(c.get('clicks') or 0 for c in campaigns)
    total_engagements = sum(c.get('engagements') or 0 for c in campaigns)
    total_spend = sum(_safe_float(c.get('cost_usd')) or 0 for c in campaigns)
    overall_ctr = calc_ctr_value(total_impressions, total_clicks)

    insights = []  # list of (text, is_positive)

    # 1. Highest CTR campaign
    campaigns_with_ctr = []
    for c in campaigns:
        ctr_val = calc_ctr_value(c.get('impressions'), c.get('clicks'))
        if ctr_val is not None:
            campaigns_with_ctr.append((c, ctr_val))
    if campaigns_with_ctr:
        best = max(campaigns_with_ctr, key=lambda x: x[1])
        best_name = best[0].get('display_name', best[0].get('name', 'Unknown'))
        camp_type = best[0].get('campaign_type', 'default')
        benchmark = CTR_BENCHMARKS.get(camp_type, CTR_BENCHMARKS['default'])
        diff = best[1] - benchmark
        above_below = 'above' if diff >= 0 else 'below'
        insights.append((
            f'{best_name} has the highest CTR at {best[1]:.2f}%, '
            f'{abs(diff):.2f}% {above_below} the benchmark',
            diff >= 0
        ))

    # 2. Total spend efficiency
    if total_spend > 0 and total_clicks > 0:
        overall_cpc = total_spend / total_clicks
        insights.append((
            f'Total spend efficiency: {format_currency(overall_cpc)} CPC across all campaigns',
            overall_cpc < 5.0  # generally good if under $5
        ))

    # 3. Impression concentration
    if len(campaigns) >= 2 and total_impressions > 0:
        top_imp_camp = max(campaigns, key=lambda c: c.get('impressions') or 0)
        top_imp = top_imp_camp.get('impressions') or 0
        pct = top_imp / total_impressions * 100
        top_name = top_imp_camp.get('display_name', top_imp_camp.get('name', 'Unknown'))
        insights.append((
            f'{pct:.0f}% of total impressions come from {top_name}',
            True
        ))

    # 4. Video completion rate
    video_campaigns = [c for c in campaigns if c.get('video_views')]
    if video_campaigns:
        total_views = sum(c.get('video_views') or 0 for c in video_campaigns)
        total_completions = sum(c.get('video_completions') or 0 for c in video_campaigns)
        if total_views > 0:
            avg_completion = total_completions / total_views * 100
            insights.append((
                f'Video completion rate averages {avg_completion:.1f}% across video campaigns',
                avg_completion >= 25.0
            ))

    # 5. Engagement rate vs industry average
    if total_impressions > 0 and total_engagements > 0:
        eng_rate = total_engagements / total_impressions * 100
        industry_avg = 0.5
        above_below = 'above' if eng_rate >= industry_avg else 'below'
        insights.append((
            f'Engagement rate of {eng_rate:.2f}% is {above_below} industry average of {industry_avg}%',
            eng_rate >= industry_avg
        ))

    # 6. Month-over-month trends (aggregate from first campaign with trends)
    for c in campaigns:
        trends = c.get('monthly_trends', [])
        if len(trends) >= 2:
            last_imp = trends[-1].get('impressions', 0)
            prev_imp = trends[-2].get('impressions', 0)
            if prev_imp > 0:
                change = (last_imp - prev_imp) / prev_imp * 100
                direction = 'grew' if change >= 0 else 'declined'
                insights.append((
                    f'Month-over-month impressions {direction} by {abs(change):.1f}%',
                    change >= 0
                ))
            break

    # Render insights (max 6)
    insights = insights[:6]
    if not insights:
        insights = [('No significant insights could be generated from the available data.', True)]

    txBox = slide.shapes.add_textbox(Emu(500000), Emu(950000), Emu(8100000), Emu(4800000))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, (text, is_positive) in enumerate(insights):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(8)
        p.space_after = Pt(8)

        # Bullet marker
        r_bullet = p.add_run()
        r_bullet.text = '  +  ' if is_positive else '  -  '
        r_bullet.font.name = FONT_NAME
        r_bullet.font.size = Pt(13)
        r_bullet.font.bold = True
        r_bullet.font.color.rgb = CLR_SUCCESS_GREEN if is_positive else CLR_DANGER_RED

        r_text = p.add_run()
        r_text.text = text
        r_text.font.name = FONT_NAME
        r_text.font.size = Pt(13)
        r_text.font.color.rgb = CLR_DARK_GRAY

    add_accent_bar(slide)
    add_logo(slide, logo_path)


def create_overall_metrics_slide(prs, campaigns, logo_path):
    """Campaign overview table with spend data."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)
    add_heading(slide, 'Campaign Overview')
    add_top_line(slide)

    has_spend = any(_safe_float(c.get('cost_usd')) for c in campaigns)

    if has_spend:
        headers = ['Campaign Name', 'Type', 'Geo', 'Status',
                   'Impressions', 'Clicks', 'CTR', 'Spend', 'CPC']
        col_widths = [Emu(1700000), Emu(950000), Emu(850000), Emu(750000),
                      Emu(950000), Emu(850000), Emu(800000), Emu(950000), Emu(840936)]
    else:
        headers = ['Campaign Name', 'Type', 'Geo', 'Start Date', 'Status',
                   'Impressions', 'Clicks', 'CTR']
        col_widths = [Emu(1900000), Emu(1050000), Emu(900000), Emu(1050000), Emu(850000),
                      Emu(1050000), Emu(1050000), Emu(990936)]

    rows = len(campaigns) + 1
    cols = len(headers)
    row_height = Emu(400000)
    table_height = min(row_height * rows, Emu(5200000))

    table_shape = slide.shapes.add_table(
        rows, cols, Emu(151532), Emu(900000), Emu(8840936), table_height
    )
    table = table_shape.table

    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        set_cell_fill(cell, CLR_LINKEDIN_BLUE)
        set_cell_text(cell, header, font_size=FONT_SIZE_SMALL, bold=True,
                     color=CLR_WHITE, alignment=PP_ALIGN.CENTER)
        set_cell_border(cell)

    for i, camp in enumerate(campaigns):
        row_idx = i + 1
        impressions = camp.get('impressions')
        clicks = camp.get('clicks')
        status = camp.get('status', 'NA')
        is_na = status in ('DRAFT', 'Yet to go Live') or impressions is None

        ctr_val = calc_ctr_value(impressions, clicks) if not is_na else None
        ctr_text = f'{ctr_val:.2f}%' if ctr_val else 'NA'
        camp_type = camp.get('campaign_type', 'Engagement')
        benchmark = CTR_BENCHMARKS.get(camp_type, CTR_BENCHMARKS['default'])
        ctr_color = CLR_BLACK
        if ctr_val is not None:
            ctr_color = CLR_SUCCESS_GREEN if ctr_val >= benchmark else CLR_DANGER_RED

        status_display = 'Active' if status == 'ACTIVE' else status.replace('_', ' ').title()

        if has_spend:
            cost = _safe_float(camp.get('cost_usd'))
            cpc = _safe_float(camp.get('cpc'))
            data = [
                camp.get('display_name', camp.get('name', 'Unknown')),
                camp_type,
                camp.get('geo_display', 'NA'),
                status_display,
                format_number(impressions) if not is_na else 'NA',
                format_number(clicks) if not is_na else 'NA',
                ctr_text,
                format_currency(cost) if cost else 'NA',
                format_currency(cpc) if cpc else 'NA',
            ]
            ctr_col = 6
        else:
            data = [
                camp.get('display_name', camp.get('name', 'Unknown')),
                camp_type,
                camp.get('geo_display', 'NA'),
                format_date_ordinal(camp.get('start_date')),
                status_display,
                format_number(impressions) if not is_na else 'NA',
                format_number(clicks) if not is_na else 'NA',
                ctr_text,
            ]
            ctr_col = 7

        for j, val in enumerate(data):
            cell = table.cell(row_idx, j)
            row_color = CLR_ROW_LIGHT if row_idx % 2 == 1 else CLR_WHITE
            set_cell_fill(cell, row_color)
            is_name = (j == 0)
            color = ctr_color if j == ctr_col else CLR_BLACK
            set_cell_text(cell, val, font_size=FONT_SIZE_SMALL,
                         bold=(is_name or j == ctr_col),
                         color=color,
                         alignment=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)
            set_cell_border(cell)

    # Apply heatmap gradient to numeric columns (impressions, clicks, CTR)
    if len(campaigns) >= 2:
        if has_spend:
            heatmap_cols = {4: 'impressions', 5: 'clicks', 6: 'ctr'}
        else:
            heatmap_cols = {5: 'impressions', 6: 'clicks', 7: 'ctr'}
        for col_idx, metric in heatmap_cols.items():
            values = []
            for camp in campaigns:
                if metric == 'ctr':
                    v = calc_ctr_value(camp.get('impressions'), camp.get('clicks')) or 0
                else:
                    v = camp.get(metric) or 0
                values.append(v)
            min_v = min(values) if values else 0
            max_v = max(values) if values else 0
            if max_v > min_v:
                for i, v in enumerate(values):
                    cell = table.cell(i + 1, col_idx)
                    set_cell_fill(cell, _heatmap_color(v, min_v, max_v))

    add_logo(slide, logo_path)


def create_impressions_chart_slide(prs, campaigns, logo_path):
    """Bar chart comparing impressions across campaigns."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)
    add_heading(slide, 'Impressions Comparison')
    add_top_line(slide)

    chart_data = CategoryChartData()
    chart_data.categories = [
        c.get('display_name', c.get('name', ''))[:20] for c in campaigns
    ]
    chart_data.add_series('Impressions', [c.get('impressions') or 0 for c in campaigns])

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Emu(500000), Emu(900000), Emu(8100000), Emu(5000000),
        chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = False
    plot = chart.plots[0]
    plot.gap_width = 80
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = CLR_LINKEDIN_BLUE
    value_axis = chart.value_axis
    value_axis.has_title = False
    value_axis.major_gridlines.format.line.color.rgb = CLR_BORDER
    category_axis = chart.category_axis
    category_axis.tick_labels.font.size = FONT_SIZE_CAPTION

    add_logo(slide, logo_path)


def create_campaign_ranking_slide(prs, campaigns, logo_path):
    """Campaign Ranking slide -- horizontal bar chart ranking by composite score."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)
    add_heading(slide, 'Campaign Ranking')
    add_top_line(slide)

    # Compute composite scores: normalize CTR, engagement rate, impressions; weight equally
    scored = []
    ctrs = []
    eng_rates = []
    imps = []

    for c in campaigns:
        ctr_val = calc_ctr_value(c.get('impressions'), c.get('clicks')) or 0.0
        er_val = _safe_float(c.get('engagement_rate')) or 0.0
        imp_val = c.get('impressions') or 0
        ctrs.append(ctr_val)
        eng_rates.append(er_val)
        imps.append(imp_val)

    max_ctr = max(ctrs) if ctrs and max(ctrs) > 0 else 1
    max_er = max(eng_rates) if eng_rates and max(eng_rates) > 0 else 1
    max_imp = max(imps) if imps and max(imps) > 0 else 1

    for i, c in enumerate(campaigns):
        norm_ctr = ctrs[i] / max_ctr
        norm_er = eng_rates[i] / max_er
        norm_imp = imps[i] / max_imp
        composite = (norm_ctr + norm_er + norm_imp) / 3.0 * 100
        name = c.get('display_name', c.get('name', 'Unknown'))[:25]
        scored.append((name, composite))

    # Sort by score descending
    scored.sort(key=lambda x: x[1], reverse=True)

    # Add rank labels to names
    ranked_names = [f'#{i+1} {s[0]}' for i, s in enumerate(scored)]
    ranked_scores = [s[1] for s in scored]

    chart_data = CategoryChartData()
    chart_data.categories = ranked_names
    chart_data.add_series('Composite Score', ranked_scores)

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Emu(400000), Emu(900000), Emu(8300000), Emu(4800000),
        chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = False
    plot = chart.plots[0]
    plot.gap_width = 60
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = CLR_LINKEDIN_BLUE
    plot.has_data_labels = True
    plot.data_labels.show_value = True
    plot.data_labels.font.size = FONT_SIZE_CAPTION
    plot.data_labels.number_format = '0.0'
    chart.value_axis.major_gridlines.format.line.color.rgb = CLR_BORDER
    chart.category_axis.tick_labels.font.size = FONT_SIZE_CAPTION

    # Subtitle
    txSub = slide.shapes.add_textbox(Emu(400000), Emu(5800000), Emu(8300000), Emu(300000))
    p = txSub.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = 'Composite score: equal-weighted average of normalized CTR, Engagement Rate, and Impressions'
    r.font.name = FONT_NAME
    r.font.size = FONT_SIZE_CAPTION
    r.font.color.rgb = CLR_SUBTLE_GRAY

    add_logo(slide, logo_path)


def create_ctr_comparison_slide(prs, campaigns, logo_path):
    """CTR Comparison slide -- grouped bar chart showing actual CTR vs benchmark."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)
    add_heading(slide, 'CTR vs Benchmark')
    add_top_line(slide)

    names = []
    actual_ctrs = []
    benchmarks = []

    for c in campaigns:
        name = c.get('display_name', c.get('name', 'Unknown'))[:20]
        ctr_val = calc_ctr_value(c.get('impressions'), c.get('clicks')) or 0.0
        camp_type = c.get('campaign_type', 'default')
        benchmark = CTR_BENCHMARKS.get(camp_type, CTR_BENCHMARKS['default'])
        names.append(name)
        actual_ctrs.append(ctr_val)
        benchmarks.append(benchmark)

    chart_data = CategoryChartData()
    chart_data.categories = names
    chart_data.add_series('Actual CTR', actual_ctrs)
    chart_data.add_series('Benchmark', benchmarks)

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Emu(500000), Emu(900000), Emu(8100000), Emu(4800000),
        chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = True
    chart.legend.include_in_layout = False
    chart.legend.font.size = FONT_SIZE_SMALL

    plot = chart.plots[0]
    plot.gap_width = 100
    plot.overlap = -20

    # Actual CTR series - LinkedIn blue
    series_actual = plot.series[0]
    series_actual.format.fill.solid()
    series_actual.format.fill.fore_color.rgb = CLR_LINKEDIN_BLUE

    # Benchmark series - light gray
    series_bench = plot.series[1]
    series_bench.format.fill.solid()
    series_bench.format.fill.fore_color.rgb = CLR_SUBTLE_GRAY

    plot.has_data_labels = True
    plot.data_labels.show_value = True
    plot.data_labels.font.size = FONT_SIZE_CAPTION
    plot.data_labels.number_format = '0.00"%"'

    chart.value_axis.major_gridlines.format.line.color.rgb = CLR_BORDER
    chart.category_axis.tick_labels.font.size = FONT_SIZE_CAPTION

    add_logo(slide, logo_path)


def create_engagement_breakdown_slide(prs, campaigns, logo_path):
    """Aggregate engagement breakdown -- pie chart + detail table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)
    add_heading(slide, 'Engagement Breakdown')
    add_top_line(slide)

    total_likes = sum(c.get('likes') or 0 for c in campaigns)
    total_comments = sum(c.get('comments') or 0 for c in campaigns)
    total_shares = sum(c.get('shares') or 0 for c in campaigns)
    total_reactions = sum(c.get('reactions') or 0 for c in campaigns)
    total_follows = sum(c.get('follows') or 0 for c in campaigns)
    total_other = sum(c.get('other_engagements') or 0 for c in campaigns)

    engagement_items = [
        ('Likes', total_likes),
        ('Comments', total_comments),
        ('Shares', total_shares),
        ('Reactions', total_reactions),
        ('Follows', total_follows),
        ('Other', total_other),
    ]
    pie_items = [(k, v) for k, v in engagement_items if v > 0]

    if pie_items:
        chart_data = CategoryChartData()
        chart_data.categories = [item[0] for item in pie_items]
        chart_data.add_series('Engagements', [item[1] for item in pie_items])

        chart_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.PIE,
            Emu(300000), Emu(850000), Emu(4500000), Emu(4500000),
            chart_data
        )
        chart = chart_frame.chart
        chart.has_legend = True
        chart.legend.include_in_layout = False
        chart.legend.font.size = FONT_SIZE_SMALL
        plot = chart.plots[0]
        plot.has_data_labels = True
        data_labels = plot.data_labels
        data_labels.show_percentage = True
        data_labels.show_value = False
        data_labels.font.size = FONT_SIZE_CAPTION

    # Detail table on the right
    total_all = sum(v for _, v in engagement_items)
    table_items = [(k, v) for k, v in engagement_items]
    rows = len(table_items) + 2  # header + data + total
    table_shape = slide.shapes.add_table(
        rows, 3, Emu(5000000), Emu(900000), Emu(3800000), Emu(300000) * rows
    )
    table = table_shape.table
    table.columns[0].width = Emu(1400000)
    table.columns[1].width = Emu(1200000)
    table.columns[2].width = Emu(1200000)

    for j, header in enumerate(['Type', 'Count', '% Share']):
        cell = table.cell(0, j)
        set_cell_fill(cell, CLR_LINKEDIN_BLUE)
        set_cell_text(cell, header, font_size=FONT_SIZE_SMALL, bold=True,
                     color=CLR_WHITE, alignment=PP_ALIGN.CENTER)
        set_cell_border(cell)

    for i, (label, count) in enumerate(table_items):
        row_idx = i + 1
        pct = f'{count / total_all * 100:.1f}%' if total_all > 0 else '0%'
        for j, val in enumerate([label, format_number(count), pct]):
            cell = table.cell(row_idx, j)
            row_color = CLR_ROW_LIGHT if row_idx % 2 == 1 else CLR_WHITE
            set_cell_fill(cell, row_color)
            set_cell_text(cell, val, font_size=FONT_SIZE_SMALL,
                         alignment=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)
            set_cell_border(cell)

    # Total row
    total_row = len(table_items) + 1
    for j, val in enumerate(['Total', format_number(total_all), '100%']):
        cell = table.cell(total_row, j)
        set_cell_fill(cell, CLR_LIGHT_GRAY)
        set_cell_text(cell, val, font_size=FONT_SIZE_SMALL, bold=True,
                     alignment=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)
        set_cell_border(cell)

    add_logo(slide, logo_path)


def create_cost_analysis_slide(prs, campaigns, logo_path):
    """Cost analysis table -- spend, CPC, CPM, cost/engagement per campaign."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)
    add_heading(slide, 'Cost & Efficiency Analysis')
    add_top_line(slide)

    headers = ['Campaign', 'Spend (USD)', 'CPC', 'CPM', 'Cost/Eng.', 'Eng. Rate']
    rows = len(campaigns) + 1
    cols = len(headers)
    row_height = Emu(400000)
    table_height = min(row_height * rows, Emu(5200000))

    table_shape = slide.shapes.add_table(
        rows, cols, Emu(200000), Emu(900000), Emu(8700000), table_height
    )
    table = table_shape.table
    col_widths = [Emu(2200000), Emu(1300000), Emu(1200000), Emu(1200000), Emu(1400000), Emu(1400000)]
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        set_cell_fill(cell, CLR_LINKEDIN_BLUE)
        set_cell_text(cell, header, font_size=FONT_SIZE_SMALL, bold=True,
                     color=CLR_WHITE, alignment=PP_ALIGN.CENTER)
        set_cell_border(cell)

    for i, camp in enumerate(campaigns):
        row_idx = i + 1
        cost = _safe_float(camp.get('cost_usd'))
        cpc = _safe_float(camp.get('cpc'))
        cpm = _safe_float(camp.get('cpm'))
        cpe = _safe_float(camp.get('cost_per_engagement'))
        eng_rate = _safe_float(camp.get('engagement_rate'))

        data = [
            camp.get('display_name', camp.get('name', 'Unknown')),
            format_currency(cost) if cost else 'NA',
            format_currency(cpc) if cpc else 'NA',
            format_currency(cpm) if cpm else 'NA',
            format_currency(cpe) if cpe else 'NA',
            f'{eng_rate:.2f}%' if eng_rate else 'NA',
        ]

        for j, val in enumerate(data):
            cell = table.cell(row_idx, j)
            row_color = CLR_ROW_LIGHT if row_idx % 2 == 1 else CLR_WHITE
            set_cell_fill(cell, row_color)
            set_cell_text(cell, val, font_size=FONT_SIZE_SMALL,
                         bold=(j == 0),
                         alignment=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)
            set_cell_border(cell)

    add_logo(slide, logo_path)


def create_campaign_title_slide(prs, campaign_name, geo, logo_path):
    """Campaign section title slide with accent bar."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)

    accent_bar = slide.shapes.add_shape(1, Emu(274320), Emu(2200000), Emu(76200), Emu(1200000))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = CLR_ORANGE
    accent_bar.line.fill.background()

    txBox = slide.shapes.add_textbox(Emu(500000), Emu(2286000), Emu(8186600), Emu(700000))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = campaign_name
    run.font.name = FONT_NAME
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = CLR_ACCENT_DARK

    txBox2 = slide.shapes.add_textbox(Emu(500000), Emu(3000000), Emu(8186600), Emu(400000))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.LEFT
    run2 = p2.add_run()
    run2.text = geo
    run2.font.name = FONT_NAME
    run2.font.size = Pt(18)
    run2.font.color.rgb = CLR_SUBTLE_GRAY

    add_accent_bar(slide)
    add_logo(slide, logo_path)


def create_performance_dashboard_slide(prs, campaign, logo_path):
    """Per-campaign performance dashboard with 8 KPI cards."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)
    add_heading(slide, 'Performance Dashboard')
    add_top_line(slide)

    imp = campaign.get('impressions')
    clicks = campaign.get('clicks')
    eng = campaign.get('engagements')
    cost = _safe_float(campaign.get('cost_usd'))
    ctr = _safe_float(campaign.get('ctr'))
    cpc = _safe_float(campaign.get('cpc'))
    cpm = _safe_float(campaign.get('cpm'))
    eng_rate = _safe_float(campaign.get('engagement_rate'))
    lpc = campaign.get('landing_page_clicks')
    conv = campaign.get('conversions')

    kpis = [
        ('Impressions', format_number(imp)),
        ('Clicks', format_number(clicks)),
        ('CTR', f'{ctr:.2f}%' if ctr else 'NA'),
        ('Engagements', format_number(eng)),
        ('Spend', format_currency(cost) if cost else 'NA'),
        ('CPC', format_currency(cpc) if cpc else 'NA'),
        ('CPM', format_currency(cpm) if cpm else 'NA'),
        ('Eng. Rate', f'{eng_rate:.2f}%' if eng_rate else 'NA'),
    ]

    card_w = Emu(2000000)
    card_h = Emu(1100000)
    gap_x = Emu(120000)
    gap_y = Emu(150000)
    start_left = Emu(300000)
    start_top = Emu(900000)

    for i, (label, value) in enumerate(kpis):
        col = i % 4
        row = i // 4
        left = start_left + col * (card_w + gap_x)
        top = start_top + row * (card_h + gap_y)
        _add_kpi_card(slide, left, top, card_w, card_h, value, label)

    # Additional metrics below cards
    extra_top = Emu(3500000)
    extras = []
    if lpc:
        extras.append(('Landing Page Clicks', format_number(lpc)))
    if conv:
        extras.append(('Conversions', format_number(conv)))
    if campaign.get('one_click_leads'):
        extras.append(('One-Click Leads', format_number(campaign.get('one_click_leads'))))
    if campaign.get('follows'):
        extras.append(('Follows', format_number(campaign.get('follows'))))

    if extras:
        txExtra = slide.shapes.add_textbox(Emu(300000), extra_top, Emu(8500000), Emu(600000))
        tf = txExtra.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        parts = []
        for label, val in extras:
            parts.append(f'{label}: {val}')
        r = p.add_run()
        r.text = '    |    '.join(parts)
        r.font.name = FONT_NAME
        r.font.size = Pt(11)
        r.font.color.rgb = CLR_DARK_GRAY

    add_accent_bar(slide)
    add_logo(slide, logo_path)


def create_general_metrics_slide(prs, campaign, sheet_url, logo_path):
    """General metrics table for a campaign + demographics link."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)
    add_heading(slide, 'General Metrics')
    add_top_line(slide)

    impressions = campaign.get('impressions')
    clicks = campaign.get('clicks')
    engagements = campaign.get('engagements')
    status = campaign.get('status', 'NA')
    is_draft = status in ('DRAFT', 'Yet to go Live') or impressions is None
    cost = _safe_float(campaign.get('cost_usd'))

    metrics = [
        ('Campaign Type', campaign.get('campaign_type', 'NA')),
        ('Geo', campaign.get('geo_display', 'NA')),
        ('Start Date', format_date_ordinal(campaign.get('start_date'))),
        ('Status', 'Active' if status == 'ACTIVE' else status.replace('_', ' ').title()),
        ('Impressions', format_number(impressions) if not is_draft else 'NA'),
        ('Clicks', format_number(clicks) if not is_draft else 'NA'),
        ('Total Engagements', format_number(engagements) if not is_draft else 'NA'),
        ('Avg. CTR', format_ctr(impressions, clicks) if not is_draft else 'NA'),
    ]
    if cost:
        metrics.append(('Total Spend', format_currency(cost)))

    rows = len(metrics) + 1
    table_shape = slide.shapes.add_table(
        rows, 2, Emu(1200000), Emu(1000000), Emu(6744000), Emu(320000) * rows
    )
    table = table_shape.table
    table.columns[0].width = Emu(3372000)
    table.columns[1].width = Emu(3372000)

    for j, header in enumerate(['Campaign Parameters', 'Results']):
        cell = table.cell(0, j)
        set_cell_fill(cell, CLR_LINKEDIN_BLUE)
        set_cell_text(cell, header, font_size=Pt(14), bold=True,
                     color=CLR_WHITE, alignment=PP_ALIGN.CENTER)
        set_cell_border(cell)

    for i, (param, value) in enumerate(metrics):
        row_idx = i + 1
        cell_param = table.cell(row_idx, 0)
        cell_value = table.cell(row_idx, 1)
        row_color = CLR_ROW_LIGHT if row_idx % 2 == 1 else CLR_WHITE
        set_cell_fill(cell_param, row_color)
        set_cell_fill(cell_value, row_color)
        set_cell_text(cell_param, param, font_size=FONT_SIZE_BODY, bold=True,
                     color=CLR_DARK_GRAY)

        val_color = CLR_BLACK
        if param == 'Avg. CTR' and not is_draft:
            ctr_val = calc_ctr_value(impressions, clicks)
            if ctr_val is not None:
                camp_type = campaign.get('campaign_type', 'Engagement')
                benchmark = CTR_BENCHMARKS.get(camp_type, CTR_BENCHMARKS['default'])
                val_color = CLR_SUCCESS_GREEN if ctr_val >= benchmark else CLR_DANGER_RED

        set_cell_text(cell_value, value, font_size=FONT_SIZE_BODY,
                     bold=(param == 'Avg. CTR'),
                     color=val_color, alignment=PP_ALIGN.CENTER)
        set_cell_border(cell_param)
        set_cell_border(cell_value)

    # Apply heatmap gradient to numeric metric cells
    heatmap_params = {'Impressions', 'Clicks', 'Total Engagements', 'Avg. CTR'}
    numeric_rows = []
    for i, (param, value) in enumerate(metrics):
        if param in heatmap_params:
            if param == 'Avg. CTR':
                num_val = calc_ctr_value(impressions, clicks) or 0
            elif param == 'Impressions':
                num_val = impressions or 0
            elif param == 'Clicks':
                num_val = clicks or 0
            elif param == 'Total Engagements':
                num_val = engagements or 0
            else:
                num_val = 0
            numeric_rows.append((i + 1, num_val))
    if numeric_rows:
        vals = [v for _, v in numeric_rows]
        min_v = min(vals)
        max_v = max(vals)
        # Only apply if there is variance (useful when viewed across context)
        if max_v > min_v:
            for row_idx, v in numeric_rows:
                cell_value = table.cell(row_idx, 1)
                set_cell_fill(cell_value, _heatmap_color(v, min_v, max_v))

    # Demographics link
    if sheet_url:
        link_top = Emu(1000000 + 320000 * rows + 200000)
        txBox2 = slide.shapes.add_textbox(Emu(1200000), link_top, Emu(3917852), Emu(369332))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        run2 = p2.add_run()
        run2.text = 'View Professional Demographics'
        run2.font.name = FONT_NAME
        run2.font.size = Pt(13)
        run2.font.bold = True
        run2.font.color.rgb = CLR_LINKEDIN_BLUE
        run2.font.underline = True
        add_hyperlink(slide, run2, sheet_url)

    add_logo(slide, logo_path)


def create_engagement_detail_slide(prs, campaign, logo_path):
    """Per-campaign engagement breakdown table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)
    add_heading(slide, 'Engagement Details')
    add_top_line(slide)

    items = [
        ('Likes', campaign.get('likes')),
        ('Comments', campaign.get('comments')),
        ('Shares', campaign.get('shares')),
        ('Reactions', campaign.get('reactions')),
        ('Follows', campaign.get('follows')),
        ('Landing Page Clicks', campaign.get('landing_page_clicks')),
        ('Company Page Clicks', campaign.get('company_page_clicks')),
        ('Other Engagements', campaign.get('other_engagements')),
        ('Conversions', campaign.get('conversions')),
        ('One-Click Leads', campaign.get('one_click_leads')),
    ]

    # Filter to items that have data
    items_with_data = [(k, v) for k, v in items if v is not None and v != 0]
    if not items_with_data:
        items_with_data = items[:6]

    rows = len(items_with_data) + 1
    table_shape = slide.shapes.add_table(
        rows, 2, Emu(2000000), Emu(1000000), Emu(5144000), Emu(320000) * rows
    )
    table = table_shape.table
    table.columns[0].width = Emu(3000000)
    table.columns[1].width = Emu(2144000)

    for j, header in enumerate(['Metric', 'Value']):
        cell = table.cell(0, j)
        set_cell_fill(cell, CLR_LINKEDIN_BLUE)
        set_cell_text(cell, header, font_size=Pt(12), bold=True,
                     color=CLR_WHITE, alignment=PP_ALIGN.CENTER)
        set_cell_border(cell)

    for i, (label, value) in enumerate(items_with_data):
        row_idx = i + 1
        cell_l = table.cell(row_idx, 0)
        cell_v = table.cell(row_idx, 1)
        row_color = CLR_ROW_LIGHT if row_idx % 2 == 1 else CLR_WHITE
        set_cell_fill(cell_l, row_color)
        set_cell_fill(cell_v, row_color)
        set_cell_text(cell_l, label, font_size=FONT_SIZE_BODY, bold=True,
                     color=CLR_DARK_GRAY)
        set_cell_text(cell_v, format_number(value), font_size=FONT_SIZE_BODY,
                     alignment=PP_ALIGN.CENTER)
        set_cell_border(cell_l)
        set_cell_border(cell_v)

    add_logo(slide, logo_path)


def create_video_performance_slide(prs, campaign, logo_path):
    """Video metrics slide -- only if campaign has video data."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)
    add_heading(slide, 'Video Performance')
    add_top_line(slide)

    views = campaign.get('video_views') or 0
    q1 = campaign.get('video_first_quartile') or 0
    mid = campaign.get('video_midpoint') or 0
    q3 = campaign.get('video_third_quartile') or 0
    completions = campaign.get('video_completions') or 0
    full_screen = campaign.get('full_screen_plays') or 0

    funnel_items = [
        ('Views', views),
        ('25% Watched', q1),
        ('50% Watched', mid),
        ('75% Watched', q3),
        ('Completed', completions),
    ]

    chart_data = CategoryChartData()
    chart_data.categories = [item[0] for item in funnel_items]
    chart_data.add_series('Video Funnel', [item[1] for item in funnel_items])

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Emu(400000), Emu(850000), Emu(5500000), Emu(4500000),
        chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = False
    plot = chart.plots[0]
    plot.gap_width = 60
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = CLR_LINKEDIN_BLUE
    plot.has_data_labels = True
    plot.data_labels.show_value = True
    plot.data_labels.font.size = FONT_SIZE_CAPTION

    # Stats on the right
    completion_rate = (completions / views * 100) if views > 0 else 0
    stats = [
        ('Total Views', format_number(views)),
        ('Completions', format_number(completions)),
        ('Completion Rate', f'{completion_rate:.1f}%'),
        ('Full Screen Plays', format_number(full_screen)),
    ]

    stats_top = Emu(1000000)
    for i, (label, value) in enumerate(stats):
        y = stats_top + i * Emu(500000)
        txLbl = slide.shapes.add_textbox(Emu(6200000), y, Emu(2700000), Emu(200000))
        p = txLbl.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = label
        r.font.name = FONT_NAME
        r.font.size = Pt(10)
        r.font.bold = True
        r.font.color.rgb = CLR_SUBTLE_GRAY

        txVal = slide.shapes.add_textbox(Emu(6200000), y + Emu(200000), Emu(2700000), Emu(250000))
        pv = txVal.text_frame.paragraphs[0]
        rv = pv.add_run()
        rv.text = value
        rv.font.name = FONT_NAME
        rv.font.size = Pt(22)
        rv.font.bold = True
        rv.font.color.rgb = CLR_ACCENT_DARK

    add_logo(slide, logo_path)


def create_monthly_trends_slide(prs, campaign, logo_path):
    """Monthly trends line chart for a campaign."""
    trends = campaign.get('monthly_trends', [])
    if not trends or len(trends) < 2:
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)
    add_heading(slide, 'Monthly Trends')
    add_top_line(slide)

    chart_data = CategoryChartData()
    chart_data.categories = [t['month'] for t in trends]
    chart_data.add_series('Impressions', [t.get('impressions', 0) for t in trends])
    chart_data.add_series('Clicks', [t.get('clicks', 0) for t in trends])

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS,
        Emu(400000), Emu(850000), Emu(8200000), Emu(4000000),
        chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = True
    chart.legend.include_in_layout = False
    chart.legend.font.size = FONT_SIZE_SMALL

    series_imp = chart.series[0]
    series_imp.format.line.color.rgb = CLR_LINKEDIN_BLUE
    series_imp.format.line.width = Emu(25400)
    series_clk = chart.series[1]
    series_clk.format.line.color.rgb = CLR_ORANGE
    series_clk.format.line.width = Emu(25400)

    chart.value_axis.major_gridlines.format.line.color.rgb = CLR_BORDER
    chart.category_axis.tick_labels.font.size = FONT_SIZE_CAPTION

    # Monthly spend trend below chart
    has_cost = any(t.get('cost', 0) > 0 for t in trends)
    if has_cost:
        txSpend = slide.shapes.add_textbox(Emu(400000), Emu(5100000), Emu(8200000), Emu(400000))
        tf = txSpend.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r1 = p.add_run()
        r1.text = 'Monthly Spend:  '
        r1.font.name = FONT_NAME
        r1.font.size = Pt(10)
        r1.font.bold = True
        r1.font.color.rgb = CLR_DARK_GRAY

        spend_parts = [f'{t["month"]}: {format_currency(t.get("cost", 0))}' for t in trends]
        r2 = p.add_run()
        r2.text = '   |   '.join(spend_parts)
        r2.font.name = FONT_NAME
        r2.font.size = Pt(9)
        r2.font.color.rgb = CLR_SUBTLE_GRAY

    # Best and worst month note
    imp_by_month = [(t.get('month', ''), t.get('impressions', 0)) for t in trends]
    if imp_by_month:
        best_month = max(imp_by_month, key=lambda x: x[1])
        worst_month = min(imp_by_month, key=lambda x: x[1])
        note_top = Emu(5500000) if has_cost else Emu(5100000)
        txNote = slide.shapes.add_textbox(Emu(400000), note_top, Emu(8200000), Emu(400000))
        tf_note = txNote.text_frame
        tf_note.word_wrap = True
        p_note = tf_note.paragraphs[0]

        r_best = p_note.add_run()
        r_best.text = f'Best month: {best_month[0]} ({format_number(best_month[1])} impressions)'
        r_best.font.name = FONT_NAME
        r_best.font.size = Pt(10)
        r_best.font.bold = True
        r_best.font.color.rgb = CLR_SUCCESS_GREEN

        r_sep = p_note.add_run()
        r_sep.text = '    |    '
        r_sep.font.name = FONT_NAME
        r_sep.font.size = Pt(10)
        r_sep.font.color.rgb = CLR_SUBTLE_GRAY

        r_worst = p_note.add_run()
        r_worst.text = f'Worst month: {worst_month[0]} ({format_number(worst_month[1])} impressions)'
        r_worst.font.name = FONT_NAME
        r_worst.font.size = Pt(10)
        r_worst.font.bold = True
        r_worst.font.color.rgb = CLR_DANGER_RED

    add_logo(slide, logo_path)


def create_creatives_slides(prs, creatives, campaign_name, logo_path):
    """Create creative ranking slides -- 2 per page."""
    if not creatives:
        return

    creatives_sorted = sorted(creatives, key=lambda c: c.get('impressions', 0), reverse=True)
    total = len(creatives_sorted)
    page_num = 0
    idx = 0

    while idx < total:
        page_num += 1
        remaining = total - idx

        if remaining >= 2:
            page_creatives = creatives_sorted[idx:idx+2]
            layout = 'two'
            idx += 2
        else:
            page_creatives = creatives_sorted[idx:idx+1]
            layout = 'one'
            idx += 1

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_slide_bg(slide)

        total_pages = math.ceil(total / 2)
        add_heading(slide, f'Creatives By Rank ({page_num}/{total_pages})')
        add_top_line(slide)
        add_logo(slide, logo_path)

        if layout == 'two' and len(page_creatives) == 2:
            positions = [
                {'left': Emu(868678), 'top': Emu(1015636), 'size': Emu(3112227)},
                {'left': Emu(5262265), 'top': Emu(1015636), 'size': Emu(3112227)}
            ]
            label_positions = [
                {'left': Emu(1500000), 'top': Emu(4389120)},
                {'left': Emu(5900000), 'top': Emu(4389120)}
            ]
        else:
            positions = [
                {'left': Emu(3016000), 'top': Emu(1015636), 'size': Emu(3112227)}
            ]
            label_positions = [
                {'left': Emu(3700000), 'top': Emu(4389120)}
            ]

        for ci, creative in enumerate(page_creatives):
            if ci >= len(positions):
                break
            pos = positions[ci]
            rank = idx - len(page_creatives) + ci + 1

            lbl_pos = label_positions[ci]
            txRank = slide.shapes.add_textbox(
                lbl_pos['left'], lbl_pos['top'], Emu(2200000), Emu(600000)
            )
            tfr = txRank.text_frame
            tfr.word_wrap = True

            pr = tfr.paragraphs[0]
            pr.alignment = PP_ALIGN.CENTER
            rr = pr.add_run()
            rr.text = f'#{rank}'
            rr.font.name = FONT_NAME
            rr.font.size = Pt(16)
            rr.font.bold = True
            rr.font.color.rgb = CLR_ACCENT_DARK

            p_stats = tfr.add_paragraph()
            p_stats.alignment = PP_ALIGN.CENTER
            r_stats = p_stats.add_run()
            imp = format_number(creative.get('impressions'))
            clk = format_number(creative.get('clicks'))
            r_stats.text = f'{imp} imp  |  {clk} clicks'
            r_stats.font.name = FONT_NAME
            r_stats.font.size = FONT_SIZE_CAPTION
            r_stats.font.color.rgb = CLR_SUBTLE_GRAY

            img_url = creative.get('image_url')
            if img_url:
                img_io = download_image(img_url)
                if img_io:
                    slide.shapes.add_picture(
                        img_io, pos['left'], pos['top'], pos['size'], pos['size']
                    )
                    # Add advertiser logo if available (#19)
                    adv_logo_url = creative.get('advertiser_logo_url')
                    if adv_logo_url:
                        logo_io = download_image(adv_logo_url)
                        if logo_io:
                            logo_size = Inches(0.5)
                            logo_left = pos['left'] + pos['size'] - logo_size
                            logo_top = pos['top']
                            slide.shapes.add_picture(
                                logo_io, logo_left, logo_top, logo_size, logo_size
                            )
                    continue

            txPlaceholder = slide.shapes.add_textbox(
                pos['left'], pos['top'] + Emu(1400000), pos['size'], Emu(400000)
            )
            tfp = txPlaceholder.text_frame
            pp = tfp.paragraphs[0]
            pp.alignment = PP_ALIGN.CENTER
            rp = pp.add_run()
            rp.text = creative.get('name', f'Creative #{rank}')
            rp.font.name = FONT_NAME
            rp.font.size = Pt(11)
            rp.font.color.rgb = CLR_DARK_GRAY


def create_top_creatives_all_campaigns_slide(prs, campaigns, logo_path):
    """Top Creatives Across All Campaigns -- top 4 by impressions with campaign label."""
    # Collect all creatives from all campaigns
    all_creatives = []
    for camp in campaigns:
        camp_name = camp.get('display_name', camp.get('name', 'Unknown'))
        for creative in camp.get('creatives', []):
            entry = dict(creative)
            entry['_campaign_name'] = camp_name
            all_creatives.append(entry)

    if not all_creatives:
        return

    # Sort by impressions, take top 4
    all_creatives.sort(key=lambda c: c.get('impressions', 0), reverse=True)
    top_creatives = all_creatives[:4]

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)
    add_heading(slide, 'Top Creatives Across All Campaigns')
    add_top_line(slide)

    count = len(top_creatives)
    if count <= 2:
        positions = [
            {'left': Emu(868678), 'top': Emu(1015636), 'size': Emu(3112227)},
            {'left': Emu(5262265), 'top': Emu(1015636), 'size': Emu(3112227)},
        ][:count]
        label_positions = [
            {'left': Emu(1000000), 'top': Emu(4389120)},
            {'left': Emu(5400000), 'top': Emu(4389120)},
        ][:count]
    else:
        # 4 in a 2x2 grid
        positions = [
            {'left': Emu(400000), 'top': Emu(950000), 'size': Emu(2200000)},
            {'left': Emu(5000000), 'top': Emu(950000), 'size': Emu(2200000)},
            {'left': Emu(400000), 'top': Emu(3500000), 'size': Emu(2200000)},
            {'left': Emu(5000000), 'top': Emu(3500000), 'size': Emu(2200000)},
        ][:count]
        label_positions = [
            {'left': Emu(2800000), 'top': Emu(950000)},
            {'left': Emu(7400000), 'top': Emu(950000)},
            {'left': Emu(2800000), 'top': Emu(3500000)},
            {'left': Emu(7400000), 'top': Emu(3500000)},
        ][:count]

    for ci, creative in enumerate(top_creatives):
        pos = positions[ci]
        lbl_pos = label_positions[ci]
        rank = ci + 1
        camp_name = creative.get('_campaign_name', '')

        # Label box
        txRank = slide.shapes.add_textbox(
            lbl_pos['left'], lbl_pos['top'], Emu(2200000), Emu(800000)
        )
        tfr = txRank.text_frame
        tfr.word_wrap = True

        pr = tfr.paragraphs[0]
        pr.alignment = PP_ALIGN.LEFT
        rr = pr.add_run()
        rr.text = f'#{rank}'
        rr.font.name = FONT_NAME
        rr.font.size = Pt(14)
        rr.font.bold = True
        rr.font.color.rgb = CLR_ACCENT_DARK

        # Campaign name
        p_camp = tfr.add_paragraph()
        p_camp.alignment = PP_ALIGN.LEFT
        rc = p_camp.add_run()
        rc.text = camp_name[:25]
        rc.font.name = FONT_NAME
        rc.font.size = Pt(9)
        rc.font.bold = True
        rc.font.color.rgb = CLR_LINKEDIN_BLUE

        # Stats
        p_stats = tfr.add_paragraph()
        p_stats.alignment = PP_ALIGN.LEFT
        r_stats = p_stats.add_run()
        imp_val = format_number(creative.get('impressions'))
        clk_val = format_number(creative.get('clicks'))
        r_stats.text = f'{imp_val} imp | {clk_val} clicks'
        r_stats.font.name = FONT_NAME
        r_stats.font.size = FONT_SIZE_CAPTION
        r_stats.font.color.rgb = CLR_SUBTLE_GRAY

        # Image
        img_url = creative.get('image_url')
        if img_url:
            img_io = download_image(img_url)
            if img_io:
                slide.shapes.add_picture(
                    img_io, pos['left'], pos['top'], pos['size'], pos['size']
                )
                continue

        # Placeholder text if no image
        txPh = slide.shapes.add_textbox(
            pos['left'], pos['top'] + Emu(900000), pos['size'], Emu(300000)
        )
        pp = txPh.text_frame.paragraphs[0]
        pp.alignment = PP_ALIGN.CENTER
        rp = pp.add_run()
        rp.text = creative.get('name', f'Creative #{rank}')
        rp.font.name = FONT_NAME
        rp.font.size = Pt(10)
        rp.font.color.rgb = CLR_DARK_GRAY

    add_logo(slide, logo_path)


def create_thankyou_slide(prs, logo_path):
    """Final Thank You slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)
    add_top_band(slide)

    txBox = slide.shapes.add_textbox(Emu(457200), Emu(2400000), Emu(8229600), Emu(914400))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = 'Thank You!'
    run.font.name = FONT_NAME
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = CLR_ACCENT_DARK

    add_orange_accent(slide, Emu(3800000), Emu(3350000), Emu(1544000))

    txSub = slide.shapes.add_textbox(Emu(457200), Emu(3500000), Emu(8229600), Emu(500000))
    tf2 = txSub.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = 'Generated with LinkedIn Report Automation'
    r2.font.name = FONT_NAME
    r2.font.size = Pt(14)
    r2.font.color.rgb = CLR_SUBTLE_GRAY

    add_accent_bar(slide)
    add_logo(slide, logo_path)


# --- HEATMAP HELPER -----------------------------------------------------------


def _heatmap_color(value, min_val, max_val):
    """Return an RGBColor on a red-yellow-green gradient.

    Lowest value = light red, middle = yellow, highest = deep green.
    """
    if max_val == min_val:
        return RGBColor(0xFF, 0xFF, 0x99)  # neutral yellow
    ratio = (value - min_val) / (max_val - min_val)  # 0..1
    if ratio <= 0.5:
        # red to yellow  (ratio 0->0.5)
        t = ratio / 0.5
        r = int(0xCC + (0xFF - 0xCC) * (1 - t))  # 0xCC..0xFF fading
        g = int(0x80 + (0xFF - 0x80) * t)
        b = int(0x80 * (1 - t))
        return RGBColor(min(r, 255), min(g, 255), min(b, 255))
    else:
        # yellow to green  (ratio 0.5->1)
        t = (ratio - 0.5) / 0.5
        r = int(0xFF * (1 - t))
        g = int(0xFF - (0xFF - 0x7A) * t)  # 0xFF -> 0x7A
        b = int(0x42 * t)
        return RGBColor(max(min(r, 255), 0), max(min(g, 255), 0), max(min(b, 255), 0))


# --- NEW SLIDE GENERATORS ----------------------------------------------------


def create_period_comparison_slide(prs, campaigns, logo_path):
    """Period-over-period comparison: current month vs previous month per campaign."""
    # Check if any campaign has monthly_trends with at least 2 entries
    has_trends = any(len(c.get('monthly_trends', [])) >= 2 for c in campaigns)
    if not has_trends:
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)
    add_heading(slide, 'Period-over-Period Comparison')
    add_top_line(slide)

    headers = ['Campaign', 'Metric', 'Previous', 'Current', '% Change']
    col_widths = [Emu(2200000), Emu(1400000), Emu(1500000), Emu(1500000), Emu(1500000)]

    # Collect rows
    row_data = []
    for camp in campaigns:
        trends = camp.get('monthly_trends', [])
        if len(trends) < 2:
            continue
        prev_month = trends[-2]
        curr_month = trends[-1]
        name = camp.get('display_name', camp.get('name', 'Unknown'))[:22]

        metrics = [
            ('Impressions', prev_month.get('impressions', 0), curr_month.get('impressions', 0)),
            ('Clicks', prev_month.get('clicks', 0), curr_month.get('clicks', 0)),
            ('CTR', prev_month.get('ctr', 0), curr_month.get('ctr', 0)),
            ('Spend', prev_month.get('cost', 0), curr_month.get('cost', 0)),
        ]
        for metric_name, prev_val, curr_val in metrics:
            prev_val = prev_val or 0
            curr_val = curr_val or 0
            if prev_val > 0:
                pct_change = (curr_val - prev_val) / prev_val * 100
            elif curr_val > 0:
                pct_change = 100.0
            else:
                pct_change = 0.0
            row_data.append((name, metric_name, prev_val, curr_val, pct_change))

    if not row_data:
        return

    rows = min(len(row_data), 16) + 1  # cap at 16 data rows + header
    cols = len(headers)
    table_height = min(Emu(350000) * rows, Emu(5400000))

    table_shape = slide.shapes.add_table(
        rows, cols, Emu(200000), Emu(900000), Emu(8100000), table_height
    )
    table = table_shape.table
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        set_cell_fill(cell, CLR_LINKEDIN_BLUE)
        set_cell_text(cell, header, font_size=FONT_SIZE_SMALL, bold=True,
                     color=CLR_WHITE, alignment=PP_ALIGN.CENTER)
        set_cell_border(cell)

    for i, (name, metric, prev_val, curr_val, pct) in enumerate(row_data[:rows - 1]):
        row_idx = i + 1
        arrow = '\u2191' if pct >= 0 else '\u2193'  # up/down arrow
        pct_color = CLR_SUCCESS_GREEN if pct >= 0 else CLR_DANGER_RED

        if metric in ('Spend',):
            prev_display = format_currency(prev_val) if prev_val else '$0.00'
            curr_display = format_currency(curr_val) if curr_val else '$0.00'
        elif metric == 'CTR':
            prev_display = f'{prev_val:.2f}%'
            curr_display = f'{curr_val:.2f}%'
        else:
            prev_display = format_number(prev_val)
            curr_display = format_number(curr_val)

        data = [name, metric, prev_display, curr_display, f'{arrow} {pct:+.1f}%']
        for j, val in enumerate(data):
            cell = table.cell(row_idx, j)
            row_color = CLR_ROW_LIGHT if row_idx % 2 == 1 else CLR_WHITE
            set_cell_fill(cell, row_color)
            color = pct_color if j == 4 else CLR_BLACK
            set_cell_text(cell, val, font_size=FONT_SIZE_SMALL,
                         bold=(j == 0 or j == 4),
                         color=color,
                         alignment=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)
            set_cell_border(cell)

    add_accent_bar(slide)
    add_logo(slide, logo_path)


def create_audience_overlap_slide(prs, campaigns, logo_path):
    """Audience overlap analysis: industries/companies appearing across multiple campaigns."""
    # Gather demographics data
    industry_map = defaultdict(set)  # industry -> set of campaign names
    company_map = defaultdict(set)   # company -> set of campaign names

    for camp in campaigns:
        demo = camp.get('demographics', {})
        name = camp.get('display_name', camp.get('name', 'Unknown'))
        for ind in demo.get('industries', []):
            ind_name = ind if isinstance(ind, str) else ind.get('name', str(ind))
            industry_map[ind_name].add(name)
        for comp in demo.get('companies', []):
            comp_name = comp if isinstance(comp, str) else comp.get('name', str(comp))
            company_map[comp_name].add(name)

    # Filter to shared audiences (appear in 2+ campaigns)
    shared_industries = {k: v for k, v in industry_map.items() if len(v) >= 2}
    shared_companies = {k: v for k, v in company_map.items() if len(v) >= 2}

    if not shared_industries and not shared_companies:
        return  # skip slide if no overlaps

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)
    add_heading(slide, 'Audience Overlap Analysis')
    add_top_line(slide)

    current_top = Emu(900000)

    # Shared Industries table
    if shared_industries:
        items = sorted(shared_industries.items(), key=lambda x: len(x[1]), reverse=True)[:8]
        rows = len(items) + 1
        table_shape = slide.shapes.add_table(
            rows, 3, Emu(300000), current_top, Emu(8400000), Emu(320000) * rows
        )
        table = table_shape.table
        table.columns[0].width = Emu(2800000)
        table.columns[1].width = Emu(4200000)
        table.columns[2].width = Emu(1400000)

        for j, header in enumerate(['Industry', 'Campaigns', '# Shared']):
            cell = table.cell(0, j)
            set_cell_fill(cell, CLR_LINKEDIN_BLUE)
            set_cell_text(cell, header, font_size=FONT_SIZE_SMALL, bold=True,
                         color=CLR_WHITE, alignment=PP_ALIGN.CENTER)
            set_cell_border(cell)

        for i, (ind_name, camp_names) in enumerate(items):
            row_idx = i + 1
            row_color = CLR_ROW_LIGHT if row_idx % 2 == 1 else CLR_WHITE
            for j, val in enumerate([ind_name, ', '.join(sorted(camp_names))[:50], str(len(camp_names))]):
                cell = table.cell(row_idx, j)
                set_cell_fill(cell, row_color)
                set_cell_text(cell, val, font_size=FONT_SIZE_SMALL,
                             alignment=PP_ALIGN.CENTER if j == 2 else PP_ALIGN.LEFT)
                set_cell_border(cell)

        current_top = current_top + Emu(320000) * rows + Emu(200000)

    # Shared Companies table
    if shared_companies and current_top < Emu(4500000):
        items = sorted(shared_companies.items(), key=lambda x: len(x[1]), reverse=True)[:6]
        rows = len(items) + 1
        table_shape = slide.shapes.add_table(
            rows, 3, Emu(300000), current_top, Emu(8400000), Emu(320000) * rows
        )
        table = table_shape.table
        table.columns[0].width = Emu(2800000)
        table.columns[1].width = Emu(4200000)
        table.columns[2].width = Emu(1400000)

        for j, header in enumerate(['Company', 'Campaigns', '# Shared']):
            cell = table.cell(0, j)
            set_cell_fill(cell, CLR_LINKEDIN_BLUE)
            set_cell_text(cell, header, font_size=FONT_SIZE_SMALL, bold=True,
                         color=CLR_WHITE, alignment=PP_ALIGN.CENTER)
            set_cell_border(cell)

        for i, (comp_name, camp_names) in enumerate(items):
            row_idx = i + 1
            row_color = CLR_ROW_LIGHT if row_idx % 2 == 1 else CLR_WHITE
            for j, val in enumerate([comp_name, ', '.join(sorted(camp_names))[:50], str(len(camp_names))]):
                cell = table.cell(row_idx, j)
                set_cell_fill(cell, row_color)
                set_cell_text(cell, val, font_size=FONT_SIZE_SMALL,
                             alignment=PP_ALIGN.CENTER if j == 2 else PP_ALIGN.LEFT)
                set_cell_border(cell)

    add_accent_bar(slide)
    add_logo(slide, logo_path)


def create_demographics_summary_slide(prs, campaign, logo_path):
    """Demographics summary: top 5 industries, job titles, companies as bar charts."""
    demo = campaign.get('demographics', {})
    if not demo:
        return

    industries = demo.get('industries', [])
    job_titles = demo.get('job_titles', [])
    companies = demo.get('companies', [])

    if not industries and not job_titles and not companies:
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)
    camp_name = campaign.get('display_name', campaign.get('name', 'Unknown'))
    add_heading(slide, f'Demographics: {camp_name[:30]}')
    add_top_line(slide)

    chart_configs = []  # (title, labels, values)

    def _extract_top5(items):
        """Extract top 5 name/value pairs from demographics list."""
        parsed = []
        for item in items:
            if isinstance(item, dict):
                name = item.get('name', str(item))[:20]
                value = item.get('count', item.get('value', item.get('percentage', 1)))
                parsed.append((name, _safe_float(value) or 1))
            elif isinstance(item, str):
                parsed.append((item[:20], 1))
        parsed.sort(key=lambda x: x[1], reverse=True)
        return parsed[:5]

    if industries:
        top5 = _extract_top5(industries)
        if top5:
            chart_configs.append(('Top Industries', [x[0] for x in top5], [x[1] for x in top5]))
    if job_titles:
        top5 = _extract_top5(job_titles)
        if top5:
            chart_configs.append(('Top Job Titles', [x[0] for x in top5], [x[1] for x in top5]))
    if companies:
        top5 = _extract_top5(companies)
        if top5:
            chart_configs.append(('Top Companies', [x[0] for x in top5], [x[1] for x in top5]))

    if not chart_configs:
        return

    num_charts = len(chart_configs)
    chart_height = Emu(1600000) if num_charts == 3 else Emu(2200000)
    chart_width = Emu(8000000)
    start_top = Emu(850000)
    gap = Emu(100000)

    for i, (title, labels, values) in enumerate(chart_configs):
        top = start_top + i * (chart_height + gap)

        chart_data = CategoryChartData()
        chart_data.categories = labels
        chart_data.add_series(title, values)

        chart_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_CLUSTERED,
            Emu(500000), top, chart_width, chart_height,
            chart_data
        )
        chart = chart_frame.chart
        chart.has_legend = False
        plot = chart.plots[0]
        plot.gap_width = 50
        series = plot.series[0]
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = CLR_LINKEDIN_BLUE
        plot.has_data_labels = True
        plot.data_labels.show_value = True
        plot.data_labels.font.size = FONT_SIZE_CAPTION
        chart.value_axis.major_gridlines.format.line.color.rgb = CLR_BORDER
        chart.category_axis.tick_labels.font.size = FONT_SIZE_CAPTION

    add_logo(slide, logo_path)


def create_campaign_funnel_slide(prs, campaign, logo_path):
    """Campaign funnel: Impressions -> Clicks -> Engagements -> Conversions."""
    imp = campaign.get('impressions') or 0
    clicks = campaign.get('clicks') or 0
    eng = campaign.get('engagements') or 0
    conv = campaign.get('conversions') or 0

    if not conv:
        return  # only show if conversion data exists

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)
    camp_name = campaign.get('display_name', campaign.get('name', 'Unknown'))
    add_heading(slide, f'Campaign Funnel: {camp_name[:28]}')
    add_top_line(slide)

    stages = [
        ('Impressions', imp),
        ('Clicks', clicks),
        ('Engagements', eng),
        ('Conversions', conv),
    ]

    chart_data = CategoryChartData()
    chart_data.categories = [s[0] for s in stages]
    chart_data.add_series('Funnel', [s[1] for s in stages])

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Emu(400000), Emu(850000), Emu(5800000), Emu(4500000),
        chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = False
    plot = chart.plots[0]
    plot.gap_width = 80
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = CLR_LINKEDIN_BLUE
    plot.has_data_labels = True
    plot.data_labels.show_value = True
    plot.data_labels.font.size = FONT_SIZE_CAPTION
    chart.value_axis.major_gridlines.format.line.color.rgb = CLR_BORDER
    chart.category_axis.tick_labels.font.size = FONT_SIZE_SMALL

    # Conversion rates on the right
    rates = []
    if imp > 0 and clicks > 0:
        rates.append(('Click Rate', clicks / imp * 100))
    if clicks > 0 and eng > 0:
        rates.append(('Engagement Rate', eng / clicks * 100))
    if eng > 0 and conv > 0:
        rates.append(('Conversion Rate', conv / eng * 100))
    if imp > 0 and conv > 0:
        rates.append(('Overall Conv.', conv / imp * 100))

    stats_top = Emu(1100000)
    for i, (label, rate) in enumerate(rates):
        y = stats_top + i * Emu(600000)
        txLbl = slide.shapes.add_textbox(Emu(6500000), y, Emu(2500000), Emu(200000))
        p = txLbl.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = label
        r.font.name = FONT_NAME
        r.font.size = Pt(10)
        r.font.bold = True
        r.font.color.rgb = CLR_SUBTLE_GRAY

        txVal = slide.shapes.add_textbox(Emu(6500000), y + Emu(200000), Emu(2500000), Emu(300000))
        pv = txVal.text_frame.paragraphs[0]
        rv = pv.add_run()
        rv.text = f'{rate:.2f}%'
        rv.font.name = FONT_NAME
        rv.font.size = Pt(20)
        rv.font.bold = True
        rv.font.color.rgb = CLR_ACCENT_DARK

    add_accent_bar(slide)
    add_logo(slide, logo_path)


def create_ai_insights_slide(prs, campaigns, logo_path):
    """AI-powered insights slide using Claude API."""
    try:
        import anthropic
        api_key = os.environ.get('ANTHROPIC_API_KEY', '')
        if not api_key:
            return
        client = anthropic.Anthropic()

        # Build metrics summary
        summary_parts = []
        for c in campaigns:
            name = c.get('display_name', c.get('name', 'Unknown'))
            imp = c.get('impressions') or 0
            clicks = c.get('clicks') or 0
            eng = c.get('engagements') or 0
            ctr = calc_ctr_value(imp, clicks)
            cost = _safe_float(c.get('cost_usd')) or 0
            eng_rate = _safe_float(c.get('engagement_rate')) or 0
            cpc = _safe_float(c.get('cpc')) or 0
            summary_parts.append(
                f"Campaign '{name}': {imp} impressions, {clicks} clicks, "
                f"CTR={ctr:.2f}% if ctr else 'N/A', {eng} engagements, "
                f"engagement_rate={eng_rate:.2f}%, spend=${cost:.2f}, CPC=${cpc:.2f}"
            )
        metrics_text = '\n'.join(summary_parts)

        message = client.messages.create(
            model="claude-sonnet-4-20250514",
            max_tokens=1024,
            messages=[{
                "role": "user",
                "content": (
                    f"You are a LinkedIn advertising analyst. Based on these campaign metrics, "
                    f"provide exactly 5 strategic insights and recommendations. "
                    f"Be specific and actionable. Keep each insight to 1-2 sentences.\n\n"
                    f"Campaign Metrics:\n{metrics_text}\n\n"
                    f"Return only the 5 insights as a numbered list (1. ... 2. ... etc.)"
                ),
            }],
        )

        response_text = message.content[0].text
        insights = []
        for line in response_text.strip().split('\n'):
            line = line.strip()
            if line and (line[0].isdigit() or line.startswith('-')):
                # Strip leading number/bullet
                clean = line.lstrip('0123456789.-) ').strip()
                if clean:
                    insights.append(clean)
        insights = insights[:5]

        if not insights:
            return

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_slide_bg(slide)
        add_heading(slide, 'AI-Powered Insights')
        add_top_line(slide)

        txBox = slide.shapes.add_textbox(Emu(500000), Emu(950000), Emu(8100000), Emu(4800000))
        tf = txBox.text_frame
        tf.word_wrap = True

        for i, insight in enumerate(insights):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.space_before = Pt(10)
            p.space_after = Pt(10)

            r_num = p.add_run()
            r_num.text = f'  {i + 1}.  '
            r_num.font.name = FONT_NAME
            r_num.font.size = Pt(13)
            r_num.font.bold = True
            r_num.font.color.rgb = CLR_LINKEDIN_BLUE

            r_text = p.add_run()
            r_text.text = insight
            r_text.font.name = FONT_NAME
            r_text.font.size = Pt(12)
            r_text.font.color.rgb = CLR_DARK_GRAY

        # Footer note
        txNote = slide.shapes.add_textbox(Emu(500000), Emu(6000000), Emu(8100000), Emu(300000))
        pn = txNote.text_frame.paragraphs[0]
        rn = pn.add_run()
        rn.text = 'Insights generated by Claude AI based on campaign performance data'
        rn.font.name = FONT_NAME
        rn.font.size = FONT_SIZE_CAPTION
        rn.font.color.rgb = CLR_SUBTLE_GRAY
        rn.font.italic = True

        add_accent_bar(slide)
        add_logo(slide, logo_path)

    except Exception:
        return  # skip silently


# --- BENCHMARKS 2026 ----------------------------------------------------------

BENCHMARKS_2026 = {
    'ctr_single_image': 0.56,
    'ctr_carousel': 0.40,
    'ctr_video': 0.44,
    'ctr_message': 3.0,
    'cpc_avg': 5.58,
    'cpm_avg': 33.80,
    'engagement_rate_non_video': 0.5,
    'engagement_rate_video': 1.6,
    'video_view_through_rate': 29.5,
    'lead_form_completion_rate': 10.0,
    'conversion_rate': 6.1,
}


# --- DEMOGRAPHIC SLIDES (PER-CAMPAIGN) --------------------------------------


def create_demo_companies_slide(prs, campaign, logo_path):
    """Top Companies by impressions for a single campaign."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)
    name = campaign.get('display_name', campaign.get('name', 'Unknown'))
    add_heading(slide, f'Top Companies \u2014 {name}')
    add_top_line(slide)

    demos = campaign.get('demographics', {})
    companies = demos.get('MEMBER_COMPANY', [])
    companies = sorted(companies, key=lambda x: x.get('impressions', 0), reverse=True)[:10]
    companies = list(reversed(companies))

    if not companies:
        txNo = slide.shapes.add_textbox(Emu(2500000), Emu(3000000), Emu(4000000), Emu(400000))
        p = txNo.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = 'No company data available'
        r.font.name = FONT_NAME
        r.font.size = Pt(14)
        r.font.color.rgb = CLR_SUBTLE_GRAY
        add_accent_bar(slide)
        add_logo(slide, logo_path)
        return

    chart_data = CategoryChartData()
    chart_data.categories = [c.get('displayName', 'Unknown')[:30] for c in companies]
    chart_data.add_series('Impressions', [c.get('impressions', 0) for c in companies])

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, Emu(400000), Emu(900000),
        Emu(8300000), Emu(5500000), chart_data
    ).chart
    chart.has_legend = False
    plot = chart.plots[0]
    plot.gap_width = 80
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = CLR_LINKEDIN_BLUE

    add_accent_bar(slide)
    add_logo(slide, logo_path)


def create_demo_industries_slide(prs, campaign, logo_path):
    """Top Industries by impressions for a single campaign."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)
    name = campaign.get('display_name', campaign.get('name', 'Unknown'))
    add_heading(slide, f'Top Industries \u2014 {name}')
    add_top_line(slide)

    demos = campaign.get('demographics', {})
    industries = demos.get('MEMBER_INDUSTRY', [])
    industries = sorted(industries, key=lambda x: x.get('impressions', 0), reverse=True)[:10]
    industries = list(reversed(industries))

    if not industries:
        txNo = slide.shapes.add_textbox(Emu(2500000), Emu(3000000), Emu(4000000), Emu(400000))
        p = txNo.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = 'No industry data available'
        r.font.name = FONT_NAME
        r.font.size = Pt(14)
        r.font.color.rgb = CLR_SUBTLE_GRAY
        add_accent_bar(slide)
        add_logo(slide, logo_path)
        return

    chart_data = CategoryChartData()
    chart_data.categories = [c.get('displayName', 'Unknown')[:30] for c in industries]
    chart_data.add_series('Impressions', [c.get('impressions', 0) for c in industries])

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, Emu(400000), Emu(900000),
        Emu(8300000), Emu(5500000), chart_data
    ).chart
    chart.has_legend = False
    plot = chart.plots[0]
    plot.gap_width = 80
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = CLR_LINKEDIN_BLUE

    add_accent_bar(slide)
    add_logo(slide, logo_path)


def create_demo_job_titles_slide(prs, campaign, logo_path):
    """Top Job Titles by impressions for a single campaign."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)
    name = campaign.get('display_name', campaign.get('name', 'Unknown'))
    add_heading(slide, f'Top Job Titles \u2014 {name}')
    add_top_line(slide)

    demos = campaign.get('demographics', {})
    titles = demos.get('MEMBER_JOB_TITLE', [])
    titles = sorted(titles, key=lambda x: x.get('impressions', 0), reverse=True)[:10]
    titles = list(reversed(titles))

    if not titles:
        txNo = slide.shapes.add_textbox(Emu(2500000), Emu(3000000), Emu(4000000), Emu(400000))
        p = txNo.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = 'No job title data available'
        r.font.name = FONT_NAME
        r.font.size = Pt(14)
        r.font.color.rgb = CLR_SUBTLE_GRAY
        add_accent_bar(slide)
        add_logo(slide, logo_path)
        return

    chart_data = CategoryChartData()
    chart_data.categories = [c.get('displayName', 'Unknown')[:30] for c in titles]
    chart_data.add_series('Impressions', [c.get('impressions', 0) for c in titles])

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, Emu(400000), Emu(900000),
        Emu(8300000), Emu(5500000), chart_data
    ).chart
    chart.has_legend = False
    plot = chart.plots[0]
    plot.gap_width = 80
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = CLR_LINKEDIN_BLUE

    add_accent_bar(slide)
    add_logo(slide, logo_path)


def create_demo_job_functions_slide(prs, campaign, logo_path):
    """Top Job Functions by impressions for a single campaign."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)
    name = campaign.get('display_name', campaign.get('name', 'Unknown'))
    add_heading(slide, f'Top Job Functions \u2014 {name}')
    add_top_line(slide)

    demos = campaign.get('demographics', {})
    functions = demos.get('MEMBER_JOB_FUNCTION', [])
    functions = sorted(functions, key=lambda x: x.get('impressions', 0), reverse=True)[:10]
    functions = list(reversed(functions))

    if not functions:
        txNo = slide.shapes.add_textbox(Emu(2500000), Emu(3000000), Emu(4000000), Emu(400000))
        p = txNo.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = 'No job function data available'
        r.font.name = FONT_NAME
        r.font.size = Pt(14)
        r.font.color.rgb = CLR_SUBTLE_GRAY
        add_accent_bar(slide)
        add_logo(slide, logo_path)
        return

    chart_data = CategoryChartData()
    chart_data.categories = [c.get('displayName', 'Unknown')[:30] for c in functions]
    chart_data.add_series('Impressions', [c.get('impressions', 0) for c in functions])

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, Emu(400000), Emu(900000),
        Emu(8300000), Emu(5500000), chart_data
    ).chart
    chart.has_legend = False
    plot = chart.plots[0]
    plot.gap_width = 80
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = CLR_LINKEDIN_BLUE

    add_accent_bar(slide)
    add_logo(slide, logo_path)


def create_demo_seniority_slide(prs, campaign, logo_path):
    """Seniority Distribution pie chart for a single campaign."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)
    name = campaign.get('display_name', campaign.get('name', 'Unknown'))
    add_heading(slide, f'Seniority Distribution \u2014 {name}')
    add_top_line(slide)

    demos = campaign.get('demographics', {})
    seniority_raw = demos.get('MEMBER_SENIORITY', [])

    # Merge duplicates (same displayName from different URNs)
    merged = {}
    for s in seniority_raw:
        dn = s.get('displayName', 'Unknown')
        merged[dn] = merged.get(dn, 0) + s.get('impressions', 0)
    seniority = [{'displayName': k, 'impressions': v} for k, v in merged.items()]
    seniority.sort(key=lambda x: x['impressions'], reverse=True)

    if not seniority:
        txNo = slide.shapes.add_textbox(Emu(2500000), Emu(3000000), Emu(4000000), Emu(400000))
        p = txNo.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = 'No seniority data available'
        r.font.name = FONT_NAME
        r.font.size = Pt(14)
        r.font.color.rgb = CLR_SUBTLE_GRAY
        add_accent_bar(slide)
        add_logo(slide, logo_path)
        return

    chart_data = CategoryChartData()
    chart_data.categories = [s.get('displayName', 'Unknown')[:30] for s in seniority]
    chart_data.add_series('Impressions', [s.get('impressions', 0) for s in seniority])

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, Emu(1200000), Emu(900000),
        Emu(6500000), Emu(5200000), chart_data
    ).chart
    chart.has_legend = True
    chart.legend.position = 2  # RIGHT
    chart.legend.include_in_layout = False
    chart.legend.font.size = FONT_SIZE_SMALL
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.show_percentage = True
    plot.data_labels.show_value = False
    plot.data_labels.font.size = FONT_SIZE_CAPTION

    add_accent_bar(slide)
    add_logo(slide, logo_path)


def create_demo_company_size_slide(prs, campaign, logo_path):
    """Company Size Distribution pie chart for a single campaign."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)
    name = campaign.get('display_name', campaign.get('name', 'Unknown'))
    add_heading(slide, f'Company Size Distribution \u2014 {name}')
    add_top_line(slide)

    demos = campaign.get('demographics', {})
    sizes_raw = demos.get('MEMBER_COMPANY_SIZE', [])

    # Merge duplicates (same displayName from different URNs)
    merged = {}
    for s in sizes_raw:
        dn = s.get('displayName', 'Unknown')
        merged[dn] = merged.get(dn, 0) + s.get('impressions', 0)
    sizes = [{'displayName': k, 'impressions': v} for k, v in merged.items()]
    sizes.sort(key=lambda x: x['impressions'], reverse=True)

    if not sizes:
        txNo = slide.shapes.add_textbox(Emu(2500000), Emu(3000000), Emu(4000000), Emu(400000))
        p = txNo.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = 'No company size data available'
        r.font.name = FONT_NAME
        r.font.size = Pt(14)
        r.font.color.rgb = CLR_SUBTLE_GRAY
        add_accent_bar(slide)
        add_logo(slide, logo_path)
        return

    chart_data = CategoryChartData()
    chart_data.categories = [s.get('displayName', 'Unknown')[:30] for s in sizes]
    chart_data.add_series('Impressions', [s.get('impressions', 0) for s in sizes])

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, Emu(1200000), Emu(900000),
        Emu(6500000), Emu(5200000), chart_data
    ).chart
    chart.has_legend = True
    chart.legend.position = 2  # RIGHT
    chart.legend.include_in_layout = False
    chart.legend.font.size = FONT_SIZE_SMALL
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.show_percentage = True
    plot.data_labels.show_value = False
    plot.data_labels.font.size = FONT_SIZE_CAPTION

    add_accent_bar(slide)
    add_logo(slide, logo_path)


def create_demo_geography_slide(prs, campaign, logo_path):
    """Geographic Distribution bar chart for a single campaign."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)
    name = campaign.get('display_name', campaign.get('name', 'Unknown'))
    add_heading(slide, f'Geographic Distribution \u2014 {name}')
    add_top_line(slide)

    demos = campaign.get('demographics', {})
    geo_data = demos.get('MEMBER_REGION_V2', demos.get('MEMBER_COUNTRY_V2', []))
    geo_data = sorted(geo_data, key=lambda x: x.get('impressions', 0), reverse=True)[:10]
    geo_data = list(reversed(geo_data))

    if not geo_data:
        txNo = slide.shapes.add_textbox(Emu(2500000), Emu(3000000), Emu(4000000), Emu(400000))
        p = txNo.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = 'No geographic data available'
        r.font.name = FONT_NAME
        r.font.size = Pt(14)
        r.font.color.rgb = CLR_SUBTLE_GRAY
        add_accent_bar(slide)
        add_logo(slide, logo_path)
        return

    chart_data = CategoryChartData()
    chart_data.categories = [c.get('displayName', 'Unknown')[:30] for c in geo_data]
    chart_data.add_series('Impressions', [c.get('impressions', 0) for c in geo_data])

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, Emu(400000), Emu(900000),
        Emu(8300000), Emu(5500000), chart_data
    ).chart
    chart.has_legend = False
    plot = chart.plots[0]
    plot.gap_width = 80
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = CLR_LINKEDIN_BLUE

    add_accent_bar(slide)
    add_logo(slide, logo_path)


def create_demo_device_slide(prs, campaign, logo_path):
    """Device Breakdown pie chart for a single campaign."""
    device_data = campaign.get('device_breakdown', [])
    if not device_data:
        return  # skip if no device data

    # Map API device names to friendly display names
    DEVICE_NAMES = {
        'MOBILE_APP': 'Mobile App',
        'DESKTOP_WEB': 'Desktop',
        'MOBILE_WEB': 'Mobile Web',
        'TABLET': 'Tablet',
        'CONNECTED_TV': 'Connected TV',
    }

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)
    name = campaign.get('display_name', campaign.get('name', 'Unknown'))
    add_heading(slide, f'Device Breakdown \u2014 {name}')
    add_top_line(slide)

    chart_data = CategoryChartData()
    labels = []
    for d in device_data:
        raw = d.get('displayName') or d.get('name') or d.get('device') or 'Unknown'
        labels.append(DEVICE_NAMES.get(raw, raw.replace('_', ' ').title())[:30])
    chart_data.categories = labels
    chart_data.add_series('Impressions', [d.get('impressions', d.get('count', 0)) for d in device_data])

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, Emu(1200000), Emu(900000),
        Emu(6500000), Emu(5200000), chart_data
    ).chart
    chart.has_legend = True
    chart.legend.position = 2  # RIGHT
    chart.legend.include_in_layout = False
    chart.legend.font.size = FONT_SIZE_SMALL
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.show_percentage = True
    plot.data_labels.show_value = False
    plot.data_labels.font.size = FONT_SIZE_CAPTION

    add_accent_bar(slide)
    add_logo(slide, logo_path)


# --- AGGREGATE DEMOGRAPHIC SLIDES -------------------------------------------


def create_aggregate_industry_slide(prs, campaigns, logo_path):
    """Aggregate top 15 industries across all campaigns."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)
    add_heading(slide, 'Top Industries \u2014 All Campaigns')
    add_top_line(slide)

    industry_totals = defaultdict(int)
    for camp in campaigns:
        demos = camp.get('demographics', {})
        for ind in demos.get('MEMBER_INDUSTRY', []):
            display = ind.get('displayName', 'Unknown')
            industry_totals[display] += ind.get('impressions', 0)

    if not industry_totals:
        txNo = slide.shapes.add_textbox(Emu(2500000), Emu(3000000), Emu(4000000), Emu(400000))
        p = txNo.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = 'No industry data available'
        r.font.name = FONT_NAME
        r.font.size = Pt(14)
        r.font.color.rgb = CLR_SUBTLE_GRAY
        add_accent_bar(slide)
        add_logo(slide, logo_path)
        return

    sorted_industries = sorted(industry_totals.items(), key=lambda x: x[1], reverse=True)[:15]
    sorted_industries = list(reversed(sorted_industries))

    chart_data = CategoryChartData()
    chart_data.categories = [ind[:30] for ind, _ in sorted_industries]
    chart_data.add_series('Impressions', [val for _, val in sorted_industries])

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, Emu(400000), Emu(900000),
        Emu(8300000), Emu(5500000), chart_data
    ).chart
    chart.has_legend = False
    plot = chart.plots[0]
    plot.gap_width = 60
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = CLR_LINKEDIN_BLUE

    add_accent_bar(slide)
    add_logo(slide, logo_path)


def create_aggregate_seniority_slide(prs, campaigns, logo_path):
    """Aggregate seniority distribution across all campaigns (pie chart)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)
    add_heading(slide, 'Seniority Distribution \u2014 All Campaigns')
    add_top_line(slide)

    seniority_totals = defaultdict(int)
    for camp in campaigns:
        demos = camp.get('demographics', {})
        for s in demos.get('MEMBER_SENIORITY', []):
            display = s.get('displayName', 'Unknown')
            seniority_totals[display] += s.get('impressions', 0)

    if not seniority_totals:
        txNo = slide.shapes.add_textbox(Emu(2500000), Emu(3000000), Emu(4000000), Emu(400000))
        p = txNo.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = 'No seniority data available'
        r.font.name = FONT_NAME
        r.font.size = Pt(14)
        r.font.color.rgb = CLR_SUBTLE_GRAY
        add_accent_bar(slide)
        add_logo(slide, logo_path)
        return

    sorted_seniority = sorted(seniority_totals.items(), key=lambda x: x[1], reverse=True)

    chart_data = CategoryChartData()
    chart_data.categories = [s[:30] for s, _ in sorted_seniority]
    chart_data.add_series('Impressions', [val for _, val in sorted_seniority])

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, Emu(1200000), Emu(900000),
        Emu(6500000), Emu(5200000), chart_data
    ).chart
    chart.has_legend = True
    chart.legend.position = 2  # RIGHT
    chart.legend.include_in_layout = False
    chart.legend.font.size = FONT_SIZE_SMALL
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.show_percentage = True
    plot.data_labels.show_value = False
    plot.data_labels.font.size = FONT_SIZE_CAPTION

    add_accent_bar(slide)
    add_logo(slide, logo_path)


def create_aggregate_company_size_slide(prs, campaigns, logo_path):
    """Aggregate company size distribution across all campaigns (pie chart)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)
    add_heading(slide, 'Company Size Distribution \u2014 All Campaigns')
    add_top_line(slide)

    size_totals = defaultdict(int)
    for camp in campaigns:
        demos = camp.get('demographics', {})
        for s in demos.get('MEMBER_COMPANY_SIZE', []):
            display = s.get('displayName', 'Unknown')
            size_totals[display] += s.get('impressions', 0)

    if not size_totals:
        txNo = slide.shapes.add_textbox(Emu(2500000), Emu(3000000), Emu(4000000), Emu(400000))
        p = txNo.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = 'No company size data available'
        r.font.name = FONT_NAME
        r.font.size = Pt(14)
        r.font.color.rgb = CLR_SUBTLE_GRAY
        add_accent_bar(slide)
        add_logo(slide, logo_path)
        return

    sorted_sizes = sorted(size_totals.items(), key=lambda x: x[1], reverse=True)

    chart_data = CategoryChartData()
    chart_data.categories = [s[:30] for s, _ in sorted_sizes]
    chart_data.add_series('Impressions', [val for _, val in sorted_sizes])

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, Emu(1200000), Emu(900000),
        Emu(6500000), Emu(5200000), chart_data
    ).chart
    chart.has_legend = True
    chart.legend.position = 2  # RIGHT
    chart.legend.include_in_layout = False
    chart.legend.font.size = FONT_SIZE_SMALL
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.show_percentage = True
    plot.data_labels.show_value = False
    plot.data_labels.font.size = FONT_SIZE_CAPTION

    add_accent_bar(slide)
    add_logo(slide, logo_path)


# --- VIRAL PERFORMANCE SLIDE ------------------------------------------------


def create_viral_performance_slide(prs, campaign, logo_path):
    """Viral Performance: organic vs viral side by side bar chart."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)
    name = campaign.get('display_name', campaign.get('name', 'Unknown'))
    add_heading(slide, f'Viral Performance \u2014 {name}')
    add_top_line(slide)

    organic_imp = (campaign.get('impressions') or 0) - (campaign.get('viral_impressions') or 0)
    organic_clicks = (campaign.get('clicks') or 0) - (campaign.get('viral_clicks') or 0)
    organic_eng = (campaign.get('engagements') or 0) - (campaign.get('viral_total_engagements') or 0)
    viral_imp = campaign.get('viral_impressions') or 0
    viral_clicks = campaign.get('viral_clicks') or 0
    viral_eng = campaign.get('viral_total_engagements') or 0

    chart_data = CategoryChartData()
    chart_data.categories = ['Impressions', 'Clicks', 'Engagements']
    chart_data.add_series('Organic', [max(organic_imp, 0), max(organic_clicks, 0), max(organic_eng, 0)])
    chart_data.add_series('Viral', [viral_imp, viral_clicks, viral_eng])

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Emu(400000), Emu(900000), Emu(5800000), Emu(4500000),
        chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = True
    chart.legend.include_in_layout = False
    chart.legend.font.size = FONT_SIZE_SMALL

    plot = chart.plots[0]
    plot.gap_width = 100
    plot.overlap = -20

    series_organic = plot.series[0]
    series_organic.format.fill.solid()
    series_organic.format.fill.fore_color.rgb = CLR_LINKEDIN_BLUE

    series_viral = plot.series[1]
    series_viral.format.fill.solid()
    series_viral.format.fill.fore_color.rgb = CLR_ORANGE

    plot.has_data_labels = True
    plot.data_labels.show_value = True
    plot.data_labels.font.size = FONT_SIZE_CAPTION

    chart.value_axis.major_gridlines.format.line.color.rgb = CLR_BORDER
    chart.category_axis.tick_labels.font.size = FONT_SIZE_SMALL

    # Viral amplification rate on the right
    viral_amp = _safe_float(campaign.get('viral_amplification_rate'))
    stats_top = Emu(1100000)
    stats = [
        ('Viral Impressions', format_number(viral_imp)),
        ('Viral Clicks', format_number(viral_clicks)),
        ('Viral Engagements', format_number(viral_eng)),
    ]
    if viral_amp is not None:
        stats.append(('Amplification Rate', f'{viral_amp:.2f}%'))

    for i, (label, value) in enumerate(stats):
        y = stats_top + i * Emu(550000)
        txLbl = slide.shapes.add_textbox(Emu(6500000), y, Emu(2500000), Emu(200000))
        p = txLbl.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = label
        r.font.name = FONT_NAME
        r.font.size = Pt(10)
        r.font.bold = True
        r.font.color.rgb = CLR_SUBTLE_GRAY

        txVal = slide.shapes.add_textbox(Emu(6500000), y + Emu(200000), Emu(2500000), Emu(250000))
        pv = txVal.text_frame.paragraphs[0]
        rv = pv.add_run()
        rv.text = value
        rv.font.name = FONT_NAME
        rv.font.size = Pt(20)
        rv.font.bold = True
        rv.font.color.rgb = CLR_ACCENT_DARK

    add_accent_bar(slide)
    add_logo(slide, logo_path)


# --- LEAD GEN FUNNEL SLIDE --------------------------------------------------


def create_lead_gen_funnel_slide(prs, campaign, logo_path):
    """Lead Generation Funnel slide with conversion rates and cost per lead."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)
    name = campaign.get('display_name', campaign.get('name', 'Unknown'))
    add_heading(slide, f'Lead Generation Funnel \u2014 {name}')
    add_top_line(slide)

    impressions = campaign.get('impressions') or 0
    lpc = campaign.get('landing_page_clicks') or 0
    form_opens = campaign.get('one_click_lead_form_opens') or 0
    leads = campaign.get('oneClickLeads', campaign.get('one_click_leads', 0)) or 0
    qualified = campaign.get('qualified_leads') or 0

    stages = [
        ('Impressions', impressions),
        ('LP Clicks', lpc),
        ('Form Opens', form_opens),
        ('Leads', leads),
    ]
    if qualified > 0:
        stages.append(('Qualified Leads', qualified))

    # Filter to stages with values
    stages = [(s, v) for s, v in stages if v > 0]
    if not stages:
        stages = [('Impressions', impressions), ('Leads', leads)]

    chart_data = CategoryChartData()
    chart_data.categories = [s for s, _ in stages]
    chart_data.add_series('Funnel', [v for _, v in stages])

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Emu(400000), Emu(900000), Emu(5800000), Emu(4500000),
        chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = False
    plot = chart.plots[0]
    plot.gap_width = 80
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = CLR_LINKEDIN_BLUE
    plot.has_data_labels = True
    plot.data_labels.show_value = True
    plot.data_labels.font.size = FONT_SIZE_CAPTION
    chart.value_axis.major_gridlines.format.line.color.rgb = CLR_BORDER
    chart.category_axis.tick_labels.font.size = FONT_SIZE_SMALL

    # Conversion rates and cost per lead on the right
    rates = []
    if impressions > 0 and lpc > 0:
        rates.append(('Click Rate', f'{lpc / impressions * 100:.2f}%'))
    if lpc > 0 and form_opens > 0:
        rates.append(('Form Open Rate', f'{form_opens / lpc * 100:.2f}%'))
    if form_opens > 0 and leads > 0:
        rates.append(('Form Completion', f'{leads / form_opens * 100:.2f}%'))
    elif impressions > 0 and leads > 0:
        rates.append(('Overall Conv.', f'{leads / impressions * 100:.4f}%'))

    cost = _safe_float(campaign.get('cost_usd'))
    if cost and leads > 0:
        cpl = cost / leads
        rates.append(('Cost per Lead', format_currency(cpl)))

    stats_top = Emu(1100000)
    for i, (label, value) in enumerate(rates):
        y = stats_top + i * Emu(550000)
        txLbl = slide.shapes.add_textbox(Emu(6500000), y, Emu(2500000), Emu(200000))
        p = txLbl.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = label
        r.font.name = FONT_NAME
        r.font.size = Pt(10)
        r.font.bold = True
        r.font.color.rgb = CLR_SUBTLE_GRAY

        txVal = slide.shapes.add_textbox(Emu(6500000), y + Emu(200000), Emu(2500000), Emu(250000))
        pv = txVal.text_frame.paragraphs[0]
        rv = pv.add_run()
        rv.text = value
        rv.font.name = FONT_NAME
        rv.font.size = Pt(20)
        rv.font.bold = True
        rv.font.color.rgb = CLR_ACCENT_DARK

    add_accent_bar(slide)
    add_logo(slide, logo_path)


# --- MESSAGING PERFORMANCE SLIDE --------------------------------------------


def create_messaging_performance_slide(prs, campaign, logo_path):
    """Messaging Ads performance funnel: Sends -> Opens -> Clicks -> Conversions."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)
    name = campaign.get('display_name', campaign.get('name', 'Unknown'))
    add_heading(slide, f'Messaging Performance \u2014 {name}')
    add_top_line(slide)

    sends = campaign.get('sends') or 0
    opens = campaign.get('opens') or 0
    clicks = campaign.get('clicks') or 0
    conv = campaign.get('conversions') or 0

    stages = [
        ('Sends', sends),
        ('Opens', opens),
        ('Clicks', clicks),
    ]
    if conv > 0:
        stages.append(('Conversions', conv))

    chart_data = CategoryChartData()
    chart_data.categories = [s for s, _ in stages]
    chart_data.add_series('Messaging Funnel', [v for _, v in stages])

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Emu(400000), Emu(900000), Emu(5800000), Emu(4500000),
        chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = False
    plot = chart.plots[0]
    plot.gap_width = 80
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = CLR_LINKEDIN_BLUE
    plot.has_data_labels = True
    plot.data_labels.show_value = True
    plot.data_labels.font.size = FONT_SIZE_CAPTION
    chart.value_axis.major_gridlines.format.line.color.rgb = CLR_BORDER
    chart.category_axis.tick_labels.font.size = FONT_SIZE_SMALL

    # Rates on the right
    rates = []
    if sends > 0 and opens > 0:
        rates.append(('Open Rate', f'{opens / sends * 100:.2f}%'))
    if opens > 0 and clicks > 0:
        rates.append(('CTR (of opens)', f'{clicks / opens * 100:.2f}%'))
    elif sends > 0 and clicks > 0:
        rates.append(('CTR (of sends)', f'{clicks / sends * 100:.2f}%'))

    stats_top = Emu(1100000)
    for i, (label, value) in enumerate(rates):
        y = stats_top + i * Emu(550000)
        txLbl = slide.shapes.add_textbox(Emu(6500000), y, Emu(2500000), Emu(200000))
        p = txLbl.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = label
        r.font.name = FONT_NAME
        r.font.size = Pt(10)
        r.font.bold = True
        r.font.color.rgb = CLR_SUBTLE_GRAY

        txVal = slide.shapes.add_textbox(Emu(6500000), y + Emu(200000), Emu(2500000), Emu(250000))
        pv = txVal.text_frame.paragraphs[0]
        rv = pv.add_run()
        rv.text = value
        rv.font.name = FONT_NAME
        rv.font.size = Pt(20)
        rv.font.bold = True
        rv.font.color.rgb = CLR_ACCENT_DARK

    add_accent_bar(slide)
    add_logo(slide, logo_path)


# --- DOCUMENT PERFORMANCE SLIDE ---------------------------------------------


def create_document_performance_slide(prs, campaign, logo_path):
    """Document Ads performance funnel: Impressions -> Q1 -> Midpoint -> Q3 -> Completions."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)
    name = campaign.get('display_name', campaign.get('name', 'Unknown'))
    add_heading(slide, f'Document Performance \u2014 {name}')
    add_top_line(slide)

    impressions = campaign.get('impressions') or 0
    q1 = campaign.get('document_first_quartile') or 0
    mid = campaign.get('document_midpoint') or 0
    q3 = campaign.get('document_third_quartile') or 0
    completions = campaign.get('document_completions') or 0

    stages = [
        ('Impressions', impressions),
        ('25% Read', q1),
        ('50% Read', mid),
        ('75% Read', q3),
        ('Completed', completions),
    ]

    chart_data = CategoryChartData()
    chart_data.categories = [s for s, _ in stages]
    chart_data.add_series('Document Funnel', [v for _, v in stages])

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Emu(400000), Emu(900000), Emu(5800000), Emu(4500000),
        chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = False
    plot = chart.plots[0]
    plot.gap_width = 80
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = CLR_LINKEDIN_BLUE
    plot.has_data_labels = True
    plot.data_labels.show_value = True
    plot.data_labels.font.size = FONT_SIZE_CAPTION
    chart.value_axis.major_gridlines.format.line.color.rgb = CLR_BORDER
    chart.category_axis.tick_labels.font.size = FONT_SIZE_SMALL

    # Completion rate on the right
    if impressions > 0:
        comp_rate = completions / impressions * 100
        txLbl = slide.shapes.add_textbox(Emu(6500000), Emu(1100000), Emu(2500000), Emu(200000))
        p = txLbl.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = 'Completion Rate'
        r.font.name = FONT_NAME
        r.font.size = Pt(10)
        r.font.bold = True
        r.font.color.rgb = CLR_SUBTLE_GRAY

        txVal = slide.shapes.add_textbox(Emu(6500000), Emu(1300000), Emu(2500000), Emu(300000))
        pv = txVal.text_frame.paragraphs[0]
        rv = pv.add_run()
        rv.text = f'{comp_rate:.1f}%'
        rv.font.name = FONT_NAME
        rv.font.size = Pt(22)
        rv.font.bold = True
        rv.font.color.rgb = CLR_ACCENT_DARK

    add_accent_bar(slide)
    add_logo(slide, logo_path)


# --- DETAILED BENCHMARKS SLIDE ----------------------------------------------


def create_detailed_benchmarks_slide(prs, campaigns, logo_path):
    """Performance vs 2026 Industry Benchmarks table with color-coding."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)
    add_heading(slide, 'Performance vs 2026 Industry Benchmarks')
    add_top_line(slide)

    total_impressions = sum(c.get('impressions') or 0 for c in campaigns)
    total_clicks = sum(c.get('clicks') or 0 for c in campaigns)
    total_engagements = sum(c.get('engagements') or 0 for c in campaigns)
    total_spend = sum(_safe_float(c.get('cost_usd')) or 0 for c in campaigns)
    total_video_views = sum(c.get('video_views') or 0 for c in campaigns)
    total_video_completions = sum(c.get('video_completions') or 0 for c in campaigns)

    overall_ctr = (total_clicks / total_impressions * 100) if total_impressions > 0 else 0
    overall_cpc = (total_spend / total_clicks) if total_clicks > 0 else 0
    overall_cpm = (total_spend / total_impressions * 1000) if total_impressions > 0 else 0
    overall_eng_rate = (total_engagements / total_impressions * 100) if total_impressions > 0 else 0
    video_vtr = (total_video_completions / total_video_views * 100) if total_video_views > 0 else 0

    # Build rows: (Metric, Your Performance, Benchmark, is_beating)
    rows_data = []

    # Overall CTR
    rows_data.append(('Overall CTR', f'{overall_ctr:.2f}%',
                      f'{BENCHMARKS_2026["ctr_single_image"]:.2f}%',
                      overall_ctr >= BENCHMARKS_2026['ctr_single_image']))

    # CPC
    if total_spend > 0 and total_clicks > 0:
        rows_data.append(('Avg. CPC', format_currency(overall_cpc),
                          format_currency(BENCHMARKS_2026['cpc_avg']),
                          overall_cpc <= BENCHMARKS_2026['cpc_avg']))  # lower is better

    # CPM
    if total_spend > 0 and total_impressions > 0:
        rows_data.append(('Avg. CPM', format_currency(overall_cpm),
                          format_currency(BENCHMARKS_2026['cpm_avg']),
                          overall_cpm <= BENCHMARKS_2026['cpm_avg']))  # lower is better

    # Engagement Rate
    if total_engagements > 0:
        bench_eng = BENCHMARKS_2026['engagement_rate_non_video']
        rows_data.append(('Engagement Rate', f'{overall_eng_rate:.2f}%',
                          f'{bench_eng:.1f}%',
                          overall_eng_rate >= bench_eng))

    # Video View Through Rate
    if total_video_views > 0:
        rows_data.append(('Video View-Through Rate', f'{video_vtr:.1f}%',
                          f'{BENCHMARKS_2026["video_view_through_rate"]:.1f}%',
                          video_vtr >= BENCHMARKS_2026['video_view_through_rate']))

    if not rows_data:
        rows_data.append(('Overall CTR', f'{overall_ctr:.2f}%',
                          f'{BENCHMARKS_2026["ctr_single_image"]:.2f}%',
                          overall_ctr >= BENCHMARKS_2026['ctr_single_image']))

    headers = ['Metric', 'Your Performance', 'Industry Benchmark', 'Status']
    rows = len(rows_data) + 1
    cols = len(headers)
    col_widths = [Emu(2500000), Emu(2000000), Emu(2200000), Emu(1500000)]
    row_height = Emu(400000)
    table_height = min(row_height * rows, Emu(5200000))

    table_shape = slide.shapes.add_table(
        rows, cols, Emu(500000), Emu(1000000), Emu(8200000), table_height
    )
    table = table_shape.table
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        set_cell_fill(cell, CLR_LINKEDIN_BLUE)
        set_cell_text(cell, header, font_size=FONT_SIZE_SMALL, bold=True,
                     color=CLR_WHITE, alignment=PP_ALIGN.CENTER)
        set_cell_border(cell)

    for i, (metric, perf, bench, beating) in enumerate(rows_data):
        row_idx = i + 1
        status_text = 'Above' if beating else 'Below'
        status_color = CLR_SUCCESS_GREEN if beating else CLR_DANGER_RED
        row_color = CLR_ROW_LIGHT if row_idx % 2 == 1 else CLR_WHITE

        for j, val in enumerate([metric, perf, bench, status_text]):
            cell = table.cell(row_idx, j)
            if j == 3:
                set_cell_fill(cell, row_color)
                set_cell_text(cell, val, font_size=FONT_SIZE_SMALL, bold=True,
                             color=status_color, alignment=PP_ALIGN.CENTER)
            else:
                set_cell_fill(cell, row_color)
                set_cell_text(cell, val, font_size=FONT_SIZE_SMALL,
                             bold=(j == 0),
                             alignment=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)
            set_cell_border(cell)

    # Footer note
    note_top = Emu(1000000) + row_height * rows + Emu(100000)
    txNote = slide.shapes.add_textbox(Emu(500000), note_top, Emu(8200000), Emu(300000))
    pn = txNote.text_frame.paragraphs[0]
    rn = pn.add_run()
    rn.text = 'Benchmarks based on 2026 LinkedIn advertising industry averages'
    rn.font.name = FONT_NAME
    rn.font.size = FONT_SIZE_CAPTION
    rn.font.color.rgb = CLR_SUBTLE_GRAY
    rn.font.italic = True

    add_accent_bar(slide)
    add_logo(slide, logo_path)


# --- MAIN ---------------------------------------------------------------------


def generate_report(input_json_path, output_pptx_path,
                    logo_path='assets/logo.png', csv_dir=None):
    """Generate the full PPTX report from JSON input."""
    with open(input_json_path, 'r', encoding='utf-8') as f:
        data = json.load(f)

    report_date = data.get('report_date', datetime.now().strftime('%d-%m-%Y'))
    campaigns = data.get('campaigns', [])
    sheet_url = data.get('google_sheet_url', '')

    # Upload demographics to Google Sheets if csv_dir is provided
    if csv_dir and os.path.isdir(csv_dir) and not sheet_url:
        try:
            from .sheets_uploader import upload_to_sheets
            sheet_url = upload_to_sheets(csv_dir, report_date)
            if sheet_url:
                data['google_sheet_url'] = sheet_url
                with open(input_json_path, 'w', encoding='utf-8') as f:
                    json.dump(data, f, indent=2)
        except Exception as e:
            print(f"Google Sheets upload failed: {e}", file=sys.stderr)
            sheet_url = ''

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # 1. Title slide
    print("  Creating title slide...", file=sys.stderr)
    create_title_slide(prs, report_date, logo_path)

    # 2. Table of Contents
    print("  Creating table of contents...", file=sys.stderr)
    create_table_of_contents_slide(prs, campaigns, logo_path)

    # 3. Executive summary
    print("  Creating executive summary...", file=sys.stderr)
    create_executive_summary_slide(prs, campaigns, logo_path)

    # 4. Key Insights
    print("  Creating key insights...", file=sys.stderr)
    create_key_insights_slide(prs, campaigns, logo_path)

    # 4b. AI-Powered Insights (if Claude API available)
    print("  Creating AI insights (if available)...", file=sys.stderr)
    create_ai_insights_slide(prs, campaigns, logo_path)

    # 5. Campaign overview table
    print("  Creating campaign overview...", file=sys.stderr)
    create_overall_metrics_slide(prs, campaigns, logo_path)

    # 6. Impressions comparison chart (2+ campaigns)
    if len(campaigns) >= 2:
        print("  Creating impressions chart...", file=sys.stderr)
        create_impressions_chart_slide(prs, campaigns, logo_path)

    # 7. Campaign Ranking (2+ campaigns)
    if len(campaigns) >= 2:
        print("  Creating campaign ranking...", file=sys.stderr)
        create_campaign_ranking_slide(prs, campaigns, logo_path)

    # 8. CTR Comparison (2+ campaigns)
    if len(campaigns) >= 2:
        print("  Creating CTR comparison...", file=sys.stderr)
        create_ctr_comparison_slide(prs, campaigns, logo_path)

    # 9. Engagement breakdown (aggregate)
    has_engagement_data = any(
        (c.get('likes') or 0) + (c.get('comments') or 0) + (c.get('shares') or 0) > 0
        for c in campaigns
    )
    if has_engagement_data:
        print("  Creating engagement breakdown...", file=sys.stderr)
        create_engagement_breakdown_slide(prs, campaigns, logo_path)

    # 10. Cost analysis (if any campaign has spend data)
    has_spend = any(_safe_float(c.get('cost_usd')) for c in campaigns)
    if has_spend:
        print("  Creating cost analysis...", file=sys.stderr)
        create_cost_analysis_slide(prs, campaigns, logo_path)

    # 10b. Period-over-period comparison
    print("  Creating period comparison...", file=sys.stderr)
    create_period_comparison_slide(prs, campaigns, logo_path)

    # 11. Per-campaign sections
    for camp in campaigns:
        name = camp.get('display_name', camp.get('name', 'Unknown'))
        geo = camp.get('geo_display', '')
        print(f"  Creating section for: {name} ({geo})...", file=sys.stderr)

        # Campaign title
        create_campaign_title_slide(prs, name, geo, logo_path)

        # Performance dashboard (KPI cards)
        create_performance_dashboard_slide(prs, camp, logo_path)

        # General metrics table
        create_general_metrics_slide(prs, camp, sheet_url, logo_path)

        # Engagement details
        has_eng = any(camp.get(k) for k in ['likes', 'comments', 'shares', 'reactions',
                                              'landing_page_clicks', 'conversions'])
        if has_eng:
            create_engagement_detail_slide(prs, camp, logo_path)

        # Campaign funnel (if conversion data exists)
        if camp.get('conversions'):
            print(f"    Adding campaign funnel for: {name}...", file=sys.stderr)
            create_campaign_funnel_slide(prs, camp, logo_path)

        # Demographics summary (if demographics data exists)
        if camp.get('demographics'):
            print(f"    Adding demographics summary for: {name}...", file=sys.stderr)
            create_demographics_summary_slide(prs, camp, logo_path)

        # Video performance (if campaign has video data)
        if camp.get('video_views'):
            print(f"    Adding video performance for: {name}...", file=sys.stderr)
            create_video_performance_slide(prs, camp, logo_path)

        # Monthly trends (if 2+ months of data)
        if len(camp.get('monthly_trends', [])) >= 2:
            print(f"    Adding monthly trends for: {name}...", file=sys.stderr)
            create_monthly_trends_slide(prs, camp, logo_path)

        # Creatives
        creatives = camp.get('creatives', [])
        if creatives:
            create_creatives_slides(prs, creatives, name, logo_path)

        # Demographics slides
        demos = camp.get('demographics', {})
        if demos:
            if demos.get('MEMBER_COMPANY'):
                print(f"    Adding top companies for: {name}...", file=sys.stderr)
                create_demo_companies_slide(prs, camp, logo_path)
            if demos.get('MEMBER_INDUSTRY'):
                print(f"    Adding top industries for: {name}...", file=sys.stderr)
                create_demo_industries_slide(prs, camp, logo_path)
            if demos.get('MEMBER_JOB_TITLE'):
                print(f"    Adding top job titles for: {name}...", file=sys.stderr)
                create_demo_job_titles_slide(prs, camp, logo_path)
            if demos.get('MEMBER_JOB_FUNCTION'):
                print(f"    Adding top job functions for: {name}...", file=sys.stderr)
                create_demo_job_functions_slide(prs, camp, logo_path)
            if demos.get('MEMBER_SENIORITY'):
                print(f"    Adding seniority distribution for: {name}...", file=sys.stderr)
                create_demo_seniority_slide(prs, camp, logo_path)
            if demos.get('MEMBER_COMPANY_SIZE'):
                print(f"    Adding company size distribution for: {name}...", file=sys.stderr)
                create_demo_company_size_slide(prs, camp, logo_path)
            if demos.get('MEMBER_REGION_V2') or demos.get('MEMBER_COUNTRY_V2'):
                print(f"    Adding geographic distribution for: {name}...", file=sys.stderr)
                create_demo_geography_slide(prs, camp, logo_path)
            if camp.get('device_breakdown'):
                print(f"    Adding device breakdown for: {name}...", file=sys.stderr)
                create_demo_device_slide(prs, camp, logo_path)

        # Viral performance (skip when viral data is negligible)
        viral_imp = camp.get('viral_impressions', 0) or 0
        if viral_imp >= 10:
            print(f"    Adding viral performance for: {name}...", file=sys.stderr)
            create_viral_performance_slide(prs, camp, logo_path)

        # Lead gen funnel
        if camp.get('one_click_lead_form_opens', 0) > 0 or camp.get('oneClickLeads', 0) > 0:
            print(f"    Adding lead gen funnel for: {name}...", file=sys.stderr)
            create_lead_gen_funnel_slide(prs, camp, logo_path)

        # Messaging performance
        if camp.get('sends', 0) > 0:
            print(f"    Adding messaging performance for: {name}...", file=sys.stderr)
            create_messaging_performance_slide(prs, camp, logo_path)

        # Document performance
        if camp.get('document_completions', 0) > 0:
            print(f"    Adding document performance for: {name}...", file=sys.stderr)
            create_document_performance_slide(prs, camp, logo_path)

    # 11b. Audience overlap analysis
    print("  Creating audience overlap analysis...", file=sys.stderr)
    create_audience_overlap_slide(prs, campaigns, logo_path)

    # Aggregate demographics
    has_demos = any(camp.get('demographics') for camp in campaigns)
    if has_demos:
        print("  Creating aggregate industry slide...", file=sys.stderr)
        create_aggregate_industry_slide(prs, campaigns, logo_path)
        print("  Creating aggregate seniority slide...", file=sys.stderr)
        create_aggregate_seniority_slide(prs, campaigns, logo_path)
        print("  Creating aggregate company size slide...", file=sys.stderr)
        create_aggregate_company_size_slide(prs, campaigns, logo_path)

    # Detailed benchmarks
    print("  Creating detailed benchmarks slide...", file=sys.stderr)
    create_detailed_benchmarks_slide(prs, campaigns, logo_path)

    # 12. Top Creatives Across All Campaigns
    has_any_creatives = any(camp.get('creatives') for camp in campaigns)
    if has_any_creatives:
        print("  Creating top creatives across all campaigns...", file=sys.stderr)
        create_top_creatives_all_campaigns_slide(prs, campaigns, logo_path)

    # 13. Thank You slide
    print("  Creating Thank You slide...", file=sys.stderr)
    create_thankyou_slide(prs, logo_path)

    # Add slide numbers to all slides except title (first) and thank you (last)
    total_slides = len(prs.slides)
    numberable_count = total_slides - 2  # exclude title and thank you
    slide_num = 0
    for i, slide in enumerate(prs.slides):
        if i == 0 or i == total_slides - 1:
            continue  # skip title and thank you
        slide_num += 1
        _add_slide_number(slide, slide_num, numberable_count)

    # Save
    prs.save(output_pptx_path)
    print(f"  Report saved to: {output_pptx_path}", file=sys.stderr)

    # Attempt PDF export
    try:
        from .pdf_exporter import export_pdf
        pdf_result = export_pdf(output_pptx_path)
        if pdf_result:
            print(f"  PDF also saved to: {pdf_result}", file=sys.stderr)
    except Exception as e:
        print(f"  PDF export skipped: {e}", file=sys.stderr)

    print(output_pptx_path)
    print(sheet_url if sheet_url else '')


if __name__ == '__main__':
    if len(sys.argv) < 3:
        print("Usage: python -m src.report_generator <input.json> <output.pptx> [logo_path] [csv_dir]",
              file=sys.stderr)
        sys.exit(1)

    input_path = sys.argv[1]
    output_path = sys.argv[2]
    logo = sys.argv[3] if len(sys.argv) > 3 else 'assets/logo.png'
    csv_dir = sys.argv[4] if len(sys.argv) > 4 else None

    generate_report(input_path, output_path, logo, csv_dir)
